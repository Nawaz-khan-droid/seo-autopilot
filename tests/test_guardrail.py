"""Guardrail enforcer tests.

Validates that the post-generation checker catches all rule violations
and that the font constants prevent regressions.
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from report.ppt_design import (
    FONT_MIN, FONT_BODY, FONT_SECTION, FONT_TITLE,
    validate_ppt_guardrail,
    add_slide_bg, _tb,
)


# ── Helpers ──

def _empty_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for _ in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_bg(slide)
    return prs


def _add_text(prs, slide_idx: int, text: str, size: int = 12):
    slide = prs.slides[slide_idx]
    _tb(slide, 0.5, 1.0, 5.0, 0.5, text, size=size)


# ── Font constant tests ──

class TestFontConstants:
    def test_min_font_gte_11(self):
        assert FONT_MIN >= 11

    def test_font_hierarchy(self):
        assert FONT_TITLE >= FONT_SECTION
        assert FONT_SECTION >= FONT_BODY
        assert FONT_BODY >= FONT_MIN

    def test_all_constants_defined(self):
        _const_map = {"FONT_MIN": FONT_MIN, "FONT_BODY": FONT_BODY, "FONT_SECTION": FONT_SECTION, "FONT_TITLE": FONT_TITLE}
        for name, val in _const_map.items():
            assert isinstance(val, int)


# ── Validator fixture tests ──

class TestValidateClean:
    """A clean PPT should pass with zero violations."""

    def test_empty_passes(self):
        prs = _empty_prs()
        violations = validate_ppt_guardrail(prs, label="Test")
        assert len(violations) == 0

    def test_normal_text_passes(self):
        prs = _empty_prs()
        _add_text(prs, 0, "Normal body text", size=FONT_BODY)
        _add_text(prs, 1, "Section Header", size=FONT_SECTION)
        _add_text(prs, 2, "Title Text", size=FONT_TITLE)
        violations = validate_ppt_guardrail(prs, label="Test")
        assert len(violations) == 0

    def test_table_with_good_font_passes(self):
        prs = _empty_prs()
        slide = prs.slides[0]
        table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(1.0), Inches(4.0), Inches(0.7))
        cell = table_shape.table.cell(0, 0)
        cell.text = "Header"
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(FONT_BODY)
        violations = validate_ppt_guardrail(prs, label="Test")
        assert len(violations) == 0


class TestValidateFontViolations:
    """Font sizes below FONT_MIN must be caught.
    
    Note: _tb() clamps sizes at call time, so violations only appear
    when fonts are set directly on paragraphs (bypassing _tb).
    """

    def test_tiny_text_caught(self):
        prs = _empty_prs()
        slide = prs.slides[0]
        from pptx.util import Pt
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5.0), Inches(0.5))
        tx.text_frame.paragraphs[0].text = "Direct tiny text"
        tx.text_frame.paragraphs[0].font.size = Pt(FONT_MIN - 1)
        violations = validate_ppt_guardrail(prs, label="Test")
        assert any("Font" in v for v in violations)

    def test_tiny_table_caught(self):
        prs = _empty_prs()
        slide = prs.slides[0]
        table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(1.0), Inches(4.0), Inches(0.7))
        cell = table_shape.table.cell(0, 0)
        cell.text = "Tiny table"
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(FONT_MIN - 2)
        violations = validate_ppt_guardrail(prs, label="Test")
        assert any("Table font" in v for v in violations)

    def test_multiple_violations_all_reported(self):
        prs = _empty_prs()
        for i in range(3):
            slide = prs.slides[min(i, 2)]
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.0 + i * 0.5), Inches(5.0), Inches(0.5))
            tx.text_frame.paragraphs[0].text = f"Direct small {i}"
            tx.text_frame.paragraphs[0].font.size = Pt(FONT_MIN - 1)
        violations = validate_ppt_guardrail(prs, label="Test")
        assert len(violations) >= 3

    def test_tb_bypass_no_false_positive(self):
        """_tb() clamps, so text added via _tb should never generate violations."""
        prs = _empty_prs()
        _add_text(prs, 0, "Clamped by tb", size=FONT_MIN - 3)
        violations = validate_ppt_guardrail(prs, label="Test")
        font_violations = [v for v in violations if "Font" in v]
        assert len(font_violations) == 0


class TestValidateOverflow:
    """'+X more' overflow pattern forbidden."""

    def test_overflow_detected(self):
        prs = _empty_prs()
        _add_text(prs, 0, "+3 more items")
        violations = validate_ppt_guardrail(prs, label="Test")
        assert any("overflow" in v.lower() for v in violations)

    def test_overflow_variants(self):
        prs = _empty_prs()
        for text in ["+ 5 more", "+12 more tasks", "+ 8 more items"]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_slide_bg(slide)
            _tb(slide, 0.5, 1.0, 5.0, 0.5, text, size=FONT_BODY)
        violations = validate_ppt_guardrail(prs, label="Test")
        overflow_violations = [v for v in violations if "overflow" in v.lower()]
        assert len(overflow_violations) >= 3


class TestValidateP0:
    """P0 priority labels are forbidden."""

    def test_p0_detected(self):
        prs = _empty_prs()
        _add_text(prs, 0, "Priority: P0")
        violations = validate_ppt_guardrail(prs, label="Test")
        assert any("P0" in v for v in violations)

    def test_p1_p2_p3_allowed(self):
        prs = _empty_prs()
        _add_text(prs, 0, "Priority: P1")
        _add_text(prs, 1, "Priority: P2")
        _add_text(prs, 2, "Priority: P3")
        violations = validate_ppt_guardrail(prs, label="Test")
        p0_violations = [v for v in violations if "P0" in v]
        assert len(p0_violations) == 0


# ── Raise-on-violation ──

class TestRaiseOnViolation:
    def test_raises_on_violation(self):
        prs = _empty_prs()
        slide = prs.slides[0]
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5.0), Inches(0.5))
        tx.text_frame.paragraphs[0].text = "Direct tiny"
        tx.text_frame.paragraphs[0].font.size = Pt(FONT_MIN - 1)
        with pytest.raises(RuntimeError, match="Guardrail failed"):
            validate_ppt_guardrail(prs, label="Test", raise_on_violation=True)

    def test_no_raise_when_clean(self):
        prs = _empty_prs()
        violations = validate_ppt_guardrail(prs, label="Test", raise_on_violation=True)
        assert len(violations) == 0


# ── Integration: _tb clamps below FONT_MIN ──

class TestTbClamp:
    def test_tb_clamps_below_min(self):
        prs = _empty_prs()
        slide = prs.slides[0]
        _tb(slide, 0.5, 1.0, 5.0, 0.5, "Should be clamped", size=FONT_MIN - 3)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if "clamped" in (p.text or ""):
                        assert p.font.size >= Pt(FONT_MIN)
                        return
        pytest.fail("Clamped text not found")
