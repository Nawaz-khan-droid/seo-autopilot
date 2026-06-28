"""Smarter report agents: builder + reviewer with LLM.

Architecture:
  ReportMaker  → builds per-section narratives from facts
  build_ppt    → assembles PPT using those narratives
  ReportReviewer → extracts text, LLM-critiques each slide, returns score
  loop (maker → build → review) until review passes or max iterations.

Guardrails built in:
  Maker: only generates content for sections with real data; skips empty sections silently
  Maker: never outputs placeholder text, pending, unavailable — enforced by system prompt + fallback
  Reviewer: rule-based pattern check for ANY placeholder, empty slides, thin content
  Reviewer: scoring floor — score<70 triggers rebuild with specific revision notes
  Both: graceful LLM failure → template fallback with real numbers only
"""

from __future__ import annotations

import io
import json
import logging
from typing import Any

from pptx import Presentation

from report.facts import ReportFacts

logger = logging.getLogger(__name__)

# ── Section labels used by both agents ──
SECTION_LABELS = [
    "executive_summary",
    "traffic_overview",
    "keyword_rankings",
    "technical_seo",
    "core_web_vitals",
    "backlink_profile",
    "competitor_snapshot",
    "next_steps",
]

# ── Placeholder patterns the reviewer flags (plus fallback checker) ──
PLACEHOLDER_PATTERNS = [
    "not available", "n/a", "pending", "no data", "connect to see",
    "data pending", "not yet available", "template", "lorem ipsum",
    "data unavailable", "insufficient data", "historical context unavailable",
    "seasonal trends", "content performance inform",
]

# ═══════════════════════════════════════════════════════════════
# 1. REPORT MAKER AGENT
# ═══════════════════════════════════════════════════════════════

MAKER_SYSTEM_PROMPT = """You are a senior SEO strategist writing slide-specific narratives for a monthly client report.

For each section you will receive raw data. Write exactly ONE concise paragraph (2-3 sentences) per section.

HARD RULES — violations will cause the report to be rejected:
1. NEVER mention "data unavailable", "insufficient data", "pending", "connect to see", "not yet available", or any variant. If data for a metric is missing, skip it silently — do not reference its absence at all.
2. NEVER use generic filler, clichés, or vague statements like "seasonal trends", "content performance", "stay competitive", "the trajectory suggests". Use ONLY the numbers and keywords in the data.
3. Every sentence must reference a specific number, keyword name, or measured metric from the data provided.
4. Use confident, direct, professional language. The client should feel informed, not sold to.
5. No markdown, no bullet points, no headings — just plain text paragraphs.
6. If the data section is empty, return an empty string.

Section-specific guidance:
- executive_summary: Overall performance, total tracked, improved/dropped counts, one key highlight.
- traffic_overview: Peak month + traffic value, month-over-month change if available, total period traffic.
- keyword_rankings: Total tracked, count improved/dropped/stable, name top mover.
- technical_seo: Health score, key issues detected, what needs attention.
- core_web_vitals: Desktop vs mobile PSI, which is stronger, any metric needing work.
- backlink_profile: Total backlinks + referring domains, dofollow/nofollow split.
- competitor_snapshot: Which competitors are tracked, any notable overlap.
- next_steps: Top priorities across teams, what should happen next month."""


def _maker_prompt(section: str, data_preview: str) -> str:
    return (
        f"Write a 2-3 sentence narrative for the '{section.replace('_', ' ').title()}' slide "
        f"of a monthly SEO client report. Use ONLY the data below. Be specific with numbers. "
        f"Skip any missing metrics without comment.\n\nData:\n{data_preview}"
    )


class ReportMakerAgent:
    """Generates per-section narrative text using LLM, with template fallback.

    Guardrails:
    - Only generates for sections with >=1 available data point
    - Empty preview → empty output (no placeholder text ever leaked)
    - LLM failure → falls back to template with real numbers only
    - Post-generation scan catches any residual placeholder text and replaces it
    """

    def __init__(self, groq_client, model: str | None = None):
        self.groq = groq_client
        self.model = model

    def run(self, facts: ReportFacts, exec_narrative: str) -> dict[str, str]:
        sections: dict[str, str] = {}
        sections["executive_summary"] = exec_narrative or self._fallback_exec(facts)

        llm_available = (self.groq is not None and hasattr(self.groq, 'chat'))

        for sec in SECTION_LABELS:
            if sec == "executive_summary":
                continue
            preview = self._preview_data(sec, facts)
            # Guardrail: skip if no real data available
            if not self._has_real_data(preview):
                sections[sec] = ""
                continue
            if llm_available:
                try:
                    prompt = _maker_prompt(sec, preview)
                    result = self.groq.chat(
                        prompt=prompt,
                        system_prompt=MAKER_SYSTEM_PROMPT,
                        max_tokens=300,
                        temperature=0.3,
                    )
                    sections[sec] = (result or "").strip()
                except Exception as e:
                    logger.warning(f"MakerAgent: {sec} LLM failed — {e}")
            if not sections.get(sec):
                sections[sec] = self._fallback_section(sec, facts)
            # Guardrail: scan output for placeholder text and sanitize
            sections[sec] = self._sanitize(sections[sec], sec, facts)

        self._log_stats(sections)
        return sections

    # ── Guardrail helpers ──

    def _has_real_data(self, preview: str) -> bool:
        """Return True if preview contains at least one numeric or named data point."""
        if not preview.strip():
            return False
        # Must contain at least one digit (number) or known metric label
        has_number = any(c.isdigit() for c in preview)
        has_label = any(kw in preview.lower() for kw in [
            "users", "clicks", "impressions", "position", "change",
            "score", "issues", "psi", "lcp", "inp", "cls",
            "backlinks", "domains", "keywords", "priority",
        ])
        return has_number or has_label

    def _sanitize(self, text: str, section: str, facts: ReportFacts) -> str:
        """Post-generation guardrail: strip any placeholder text from output."""
        if not text:
            return ""
        lower = text.lower()
        for pat in PLACEHOLDER_PATTERNS:
            if pat in lower:
                logger.warning(f"MakerAgent: {section} contained '{pat}' — replacing with fallback")
                return self._fallback_section(section, facts)
        return text

    # ── data preview builders ──

    def _preview_data(self, section: str, facts: ReportFacts) -> str:
        lines: list[str] = []
        if section == "traffic_overview":
            mt = facts.kpis.monthly_traffic
            if mt:
                for p in mt:
                    if p.users.is_available:
                        lines.append(f"{p.month}: {p.users.value} users")
            if facts.kpis.clicks.is_available:
                lines.append(f"Clicks: {facts.kpis.clicks.value}")
            if facts.kpis.impressions.is_available:
                lines.append(f"Impressions: {facts.kpis.impressions.value}")
        elif section == "keyword_rankings":
            lines.append(f"Total tracked: {len(facts.rankings)}")
            for r in facts.rankings[:8]:
                p = r.position.value if r.position.is_available else "?"
                c = r.change.value if r.change.is_available else "0"
                lines.append(f"{r.keyword}: pos {p}, change {c}")
        elif section == "technical_seo":
            t = facts.technical
            if t.health_score.is_available:
                lines.append(f"Health score: {t.health_score.value}/100")
            if t.total_issues.is_available:
                lines.append(f"Total issues: {t.total_issues.value}")
            if t.missing_h1.is_available:
                lines.append(f"Missing H1: {t.missing_h1.value}")
            if t.missing_meta.is_available:
                lines.append(f"Missing meta: {t.missing_meta.value}")
        elif section == "core_web_vitals":
            c = facts.cwv
            if c.desktop_score.is_available:
                lines.append(f"Desktop PSI: {c.desktop_score.value}/100")
            if c.mobile_score.is_available:
                lines.append(f"Mobile PSI: {c.mobile_score.value}/100")
            if c.lcp_seconds.is_available:
                lines.append(f"LCP: {c.lcp_seconds.value}s")
            if c.inp_ms.is_available:
                lines.append(f"INP: {c.inp_ms.value}ms")
            if c.cls_score.is_available:
                lines.append(f"CLS: {c.cls_score.value}")
        elif section == "backlink_profile":
            ba = facts.backlinks
            if ba.total_backlinks.is_available:
                lines.append(f"Total backlinks: {ba.total_backlinks.value}")
            if ba.ref_domains.is_available:
                lines.append(f"Ref domains: {ba.ref_domains.value}")
            if ba.dofollow_count.is_available:
                lines.append(f"DoFollow: {ba.dofollow_count.value}")
            if ba.nofollow_count.is_available:
                lines.append(f"NoFollow: {ba.nofollow_count.value}")
        elif section == "competitor_snapshot":
            # Aggregate unique competitors from all rankings
            seen = set()
            for r in facts.rankings:
                for c in r.competitors:
                    if c and c not in seen:
                        seen.add(c)
                        lines.append(c)
            if not lines:
                lines.append("No competitor data available")
        elif section == "next_steps":
            for item in facts.action_plan:
                # Include team, priority, impact, effort, owner, eta
                parts = []
                if item.team:
                    parts.append(f"[{item.team}]")
                if item.priority:
                    parts.append(f"P{item.priority}")
                if item.impact:
                    parts.append(f"Impact:{item.impact}")
                if item.effort:
                    parts.append(f"Effort:{item.effort}")
                if item.owner and item.owner != "Unassigned":
                    parts.append(f"Owner:{item.owner}")
                if item.eta and item.eta != "TBD":
                    parts.append(f"ETA:{item.eta}")
                task = item.task or "Untitled task"
                lines.append(f"{' '.join(parts)} {task}")
        return "\n".join(lines)

    # ── Fallback generators (used when LLM unavailable or output rejected) ──

    def _fallback_exec(self, facts: ReportFacts) -> str:
        imp = sum(1 for r in facts.rankings if r.change.is_available and str(r.change.value).startswith("+"))
        drp = sum(1 for r in facts.rankings if r.change.is_available and str(r.change.value).startswith("-"))
        nf = sum(1 for r in facts.rankings if str(r.position.value).strip().lower() == "not found")
        return (
            f"We tracked {len(facts.rankings)} keywords this period. "
            f"{imp} improved, {drp} dropped, and {nf} were not found. "
            "See the detailed sections below for per-keyword analysis and recommendations."
        )

    def _fallback_section(self, section: str, facts: ReportFacts) -> str:
        if section == "traffic_overview":
            mt = facts.kpis.monthly_traffic
            if mt and any(p.users.is_available for p in mt):
                peak = max(
                    (p for p in mt if p.users.is_available),
                    key=lambda p: int(str(p.users.value).replace(",", "") or 0),
                    default=None,
                )
                if peak:
                    return f"Peak traffic reached {peak.users.value} users in {peak.month}."
            return ""
        if section == "keyword_rankings":
            imp = sum(1 for r in facts.rankings if r.change.is_available and str(r.change.value).startswith("+"))
            drp = sum(1 for r in facts.rankings if r.change.is_available and str(r.change.value).startswith("-"))
            stable = len(facts.rankings) - imp - drp
            return f"Of {len(facts.rankings)} tracked keywords, {imp} gained positions, {drp} lost, and {stable} held steady."
        if section == "technical_seo":
            t = facts.technical
            if t.health_score.is_available:
                return f"Site health score is {t.health_score.value}/100 with {t.total_issues.value or 'N/A'} issues detected."
            return ""
        if section == "core_web_vitals":
            c = facts.cwv
            parts = []
            if c.desktop_score.is_available:
                parts.append(f"Desktop PSI {c.desktop_score.value}/100")
            if c.mobile_score.is_available:
                parts.append(f"Mobile PSI {c.mobile_score.value}/100")
            for label, val, unit in [
                ("LCP", c.lcp_seconds, "s"),
                ("INP", c.inp_ms, "ms"),
                ("CLS", c.cls_score, ""),
            ]:
                if val.is_available:
                    parts.append(f"{label} {val.value}{unit}")
            return " | ".join(parts) if parts else ""
        if section == "backlink_profile":
            ba = facts.backlinks
            parts = []
            if ba.total_backlinks.is_available:
                parts.append(f"{ba.total_backlinks.value} total")
            if ba.ref_domains.is_available:
                parts.append(f"{ba.ref_domains.value} ref domains")
            if ba.dofollow_count.is_available:
                parts.append(f"{ba.dofollow_count.value} dofollow")
            return " | ".join(parts) if parts else ""
        if section == "competitor_snapshot":
            seen = set()
            for r in facts.rankings:
                for c in r.competitors:
                    if c and c not in seen:
                        seen.add(c)
            if seen:
                return f"Tracking {len(seen)} competitor domains across keywords."
            return ""
        if section == "next_steps":
            items = []
            for a in facts.action_plan[:3]:
                parts = []
                if a.team:
                    parts.append(f"[{a.team}]")
                if a.priority:
                    parts.append(f"P{a.priority}")
                if a.impact:
                    parts.append(f"Impact:{a.impact}")
                if a.effort:
                    parts.append(f"Effort:{a.effort}")
                task = a.task or "Untitled"
                items.append(f"{' '.join(parts)} {task}")
            return " ".join(f"{i+1}. {t}" for i, t in enumerate(items)) if items else ""
        return ""

    def _log_stats(self, sections: dict[str, str]) -> None:
        for sec, txt in sections.items():
            logger.info(f"MakerAgent: {sec} = {len(txt)} chars | {txt[:100]!r}")


# ═══════════════════════════════════════════════════════════════
# 2. REPORT REVIEWER AGENT
# ═══════════════════════════════════════════════════════════════

REVIEWER_SYSTEM_PROMPT = """You are a strict quality assurance reviewer for monthly SEO client reports.

You will receive the extracted text content from each slide of a generated PowerPoint report.
Review each slide against the checklist below. Return a JSON response with exactly these keys:

{
  "score": <integer 0-100>,
  "approved": <true/false>,
  "issues": ["Issue on Slide X: description", ...],
  "revision_notes": "<concise revision guidance if not approved, else ''>"
}

REJECT CONDITIONS (score < 70 = not approved):
1. Any slide contains: "not available", "n/a", "pending", "no data", "connect to see", "data unavailable", "insufficient data", "historical context", "seasonal trends", or any placeholder language.
2. Any slide has only a heading with no body content, or content is just 1-2 words.
3. Executive summary slide lacks specific numbers (total tracked, improved, dropped counts).
4. Multiple slides rely on generic filler instead of specific data.
5. Narrative reads like marketing copy instead of a data-driven analysis.
6. Numbers in the narrative are vague ("many", "several", "some") instead of specific.

SCORING:
- 90-100: All slides data-driven, specific, professional. No issues.
- 70-89: Minor issues (vague phrasing, missing number on 1 slide). Approve with notes.
- 50-69: Significant issues (placeholders, thin content on multiple slides). Reject.
- <50: Critical issues. Report needs full rewrite.

Return ONLY the JSON object with no surrounding text."""


class ReportReviewerAgent:
    """Reviews generated PPTX content using LLM, returns score + issues.

    Guardrails:
    - LLM review with strict checklist
    - Rule-based fallback scans every slide for placeholder patterns, empty content, thin content
    - Score < 70 → not approved → triggers rebuild
    """

    def __init__(self, groq_client, model: str | None = None):
        self.groq = groq_client
        self.model = model

    def review(self, ppt_bytes: bytes) -> dict[str, Any]:
        try:
            prs = Presentation(io.BytesIO(ppt_bytes))
        except Exception as e:
            return {
                "score": 0, "approved": False,
                "issues": [f"Cannot open PPTX: {e}"],
                "revision_notes": "Invalid PPTX file — cannot review.",
            }

        slide_texts: list[dict[str, Any]] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = (p.text or "").strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = (cell.text or "").strip()
                            if t:
                                texts.append(t)
            slide_texts.append({"slide": i, "texts": texts})

        # Rule-based guardrail check always runs
        rule_issues = self._rule_check(slide_texts)

        # LLM review
        review_data = self._build_review_payload(slide_texts)
        result = self._llm_review(review_data)

        if result:
            # Merge rule issues into result
            for iss in rule_issues:
                if iss not in result.get("issues", []):
                    result.setdefault("issues", []).append(iss)
            # Recalculate score based on issues
            result["score"] = self._compute_score(result.get("issues", []))
            result["approved"] = result["score"] >= 70
        else:
            # LLM unavailable: use rule-based result as primary
            score = self._compute_score(rule_issues)
            result = {
                "score": score,
                "approved": score >= 70,
                "issues": rule_issues,
                "revision_notes": "" if score >= 70
                    else f"Found {len(rule_issues)} issues. Review slides for placeholders or thin content.",
            }

        logger.info(
            f"ReviewerAgent: score={result.get('score')} "
            f"approved={result.get('approved')} "
            f"issues={len(result.get('issues', []))}"
        )
        return result

    # ── Rule-based guardrail ──

    def _rule_check(self, slide_texts: list[dict]) -> list[str]:
        issues: list[str] = []
        for s in slide_texts:
            num = s["slide"]
            all_text = " ".join(s["texts"]).lower()
            # Check placeholder patterns
            for pat in PLACEHOLDER_PATTERNS:
                if pat in all_text:
                    issues.append(f"Slide {num}: contains placeholder text '{pat}'")
                    break
            # Check for empty slides (header only or no content)
            content_lines = [t for t in s["texts"] if len(t) > 3]
            if len(content_lines) <= 1:
                issues.append(f"Slide {num}: only 1 content line — may be empty or thin")
            # Check for slides with no numeric data
            has_digit = any(c.isdigit() for c in all_text)
            if not has_digit and num > 1:  # skip cover slide
                issues.append(f"Slide {num}: no numeric data found — may lack specific numbers")
        return issues

    def _compute_score(self, issues: list[str]) -> int:
        base = 100
        for _ in issues:
            base -= 12
        return max(0, min(100, base))

    # ── LLM review ──

    def _build_review_payload(self, slide_texts: list[dict]) -> str:
        lines = [f"Total slides: {len(slide_texts)}\n"]
        for s in slide_texts:
            lines.append(f"--- Slide {s['slide']} ---")
            for t in s["texts"][:25]:
                lines.append(f"  {t[:150]}")
            lines.append("")
        return "\n".join(lines)

    def _llm_review(self, payload: str) -> dict[str, Any] | None:
        if self.groq is None or not hasattr(self.groq, 'chat'):
            return None
        prompt = (
            f"Review this SEO report slide content against the checklist. "
            f"Return JSON with 'score', 'approved', 'issues', 'revision_notes'.\n\n"
            f"Slide content:\n{payload}"
        )
        try:
            raw = self.groq.chat(
                prompt=prompt,
                system_prompt=REVIEWER_SYSTEM_PROMPT,
                max_tokens=800,
                temperature=0.1,
            )
            if raw:
                start = raw.index("{")
                end = raw.rindex("}") + 1
                parsed = json.loads(raw[start:end])
                if "score" in parsed and "approved" in parsed:
                    return parsed
        except Exception as e:
            logger.warning(f"ReviewerAgent: LLM parse failed — {e}")
        return None
