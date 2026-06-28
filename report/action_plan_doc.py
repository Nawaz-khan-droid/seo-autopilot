"""Monthly Action Plan — Internal team briefing, 6-10 slides.

Voice: Smart Strategy Advisor / Senior Analyst guiding the internal team.
Each task, dependency, and recommendation originates from ReportFacts ActionItems.
No hallucinated claims. No placeholder language.
"""
from __future__ import annotations

import io
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from report.facts import ActionItem, ReportFacts
from report.ppt_design import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_W, CONTENT_L,
    NAVY, TEAL, AMBER, CORAL, WHITE, GRAY_BG, GRAY_LINE, PRIMARY_BLUE,
    GRAY_MID, GRAY_TEXT, DARK_TEXT, INFO_BG,
    add_slide_bg, add_header_bar, add_footer, _tb, _rect, _line,
)

logger = logging.getLogger(__name__)

PRIO_SORT = {"P0": 0, "P1": 1, "P2": 2, "P3": 3}

NAVY_LIGHT = RGBColor(0x2C, 0x3E, 0x66)


def _priority_color(prio: str) -> RGBColor:
    return {"P1": CORAL, "P2": AMBER, "P3": TEAL}.get(prio, GRAY_MID)


def _impact_color(impact: str) -> RGBColor:
    return {"high": CORAL, "medium": AMBER, "low": TEAL}.get(impact.lower(), GRAY_MID)


def _advisor_note(slide, l: float, t: float, w: float, h: float,
                  heading: str, body: str, accent=PRIMARY_BLUE):
    """Smart advisor callout — gives context, warnings, or suggestions."""
    _rect(slide, l, t, w, h, INFO_BG)
    _rect(slide, l, t, 0.06, h, accent)
    _tb(slide, l + 0.18, t + 0.08, w - 0.3, 0.22,
        heading, size=10, bold=True, color=accent)
    _tb(slide, l + 0.18, t + 0.32, w - 0.3, h - 0.4,
        body, size=9, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def _cover_slide(prs: Presentation, client: str, month: str, agency: str) -> int:
    """Cover: internal team briefing header."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, NAVY)
    _rect(slide, 0, 0, SLIDE_W, 0.06, TEAL)
    _tb(slide, MARGIN, 1.5, CONTENT_W, 0.35, "MONTHLY ACTION PLAN",
        size=14, bold=True, color=TEAL)
    _tb(slide, MARGIN, 2.0, CONTENT_W, 0.8, client,
        size=36, bold=True, color=WHITE)
    _tb(slide, MARGIN, 2.8, CONTENT_W, 0.4, month,
        size=16, color=RGBColor(0xB0, 0xC0, 0xD0))
    _rect(slide, MARGIN, 3.4, 2.5, 0.03, TEAL)
    _tb(slide, MARGIN, 3.7, CONTENT_W, 0.4,
        "Strategy Advisor Briefing — Internal Use",
        size=11, color=GRAY_MID)
    _rect(slide, 0, SLIDE_H - 0.4, SLIDE_W, 0.4, NAVY_LIGHT)
    return 1


def _overview_slide(prs: Presentation, items: list[ActionItem],
                    page: int) -> int:
    """Slide 2: Task landscape — what the team needs to deliver."""
    page += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_header_bar(slide, "Task Landscape",
                   f"{len(items)} action items this period — breakdown by team and priority",
                   page, "Overview")

    # Team breakdown
    teams: dict[str, int] = {}
    prios: dict[str, int] = {"P1": 0, "P2": 0, "P3": 0}
    for item in items:
        teams[item.team] = teams.get(item.team, 0) + 1
        prio = item.priority if item.priority in prios else "P3"
        prios[prio] = prios.get(prio, 0) + 1

    y = 1.3
    col_w = (CONTENT_W - 0.5) / 2

    _tb(slide, MARGIN, y, col_w, 0.3, "By Team", size=14, bold=True, color=NAVY)
    y += 0.4
    for team, count in sorted(teams.items(), key=lambda x: -x[1]):
        team_color = {"SEO": TEAL, "Dev": CORAL, "Content": AMBER, "Local": NAVY}.get(team, GRAY_MID)
        _rect(slide, MARGIN + 0.1, y, 0.06, 0.35, team_color)
        _tb(slide, MARGIN + 0.25, y + 0.04, 2.0, 0.25,
            team, size=11, bold=True, color=DARK_TEXT)
        _tb(slide, MARGIN + 2.5, y + 0.04, 0.6, 0.25,
            str(count), size=14, bold=True, color=team_color)
        y += 0.4

    rx = MARGIN + col_w + 0.5
    y2 = 1.3
    _tb(slide, rx, y2, col_w, 0.3, "By Urgency", size=14, bold=True, color=NAVY)
    y2 += 0.4
    for prio in ["P1", "P2", "P3"]:
        count = prios.get(prio, 0)
        pc = _priority_color(prio)
        _rect(slide, rx + 0.1, y2, 0.06, 0.35, pc)
        _tb(slide, rx + 0.25, y2 + 0.04, 2.0, 0.25,
            prio, size=11, bold=True, color=DARK_TEXT)
        _tb(slide, rx + 2.5, y2 + 0.04, 0.6, 0.25,
            str(count), size=14, bold=True, color=pc)
        y2 += 0.4

    # Advisor analysis
    y3 = max(y, y2) + 0.3
    _line(slide, MARGIN, y3, CONTENT_W)

    p1_high = sum(1 for a in items if a.priority == "P1" and a.impact == "high")
    advisory_lines = []
    if p1_high > 0:
        advisory_lines.append(
            f"Urgent: {p1_high} P1-high items need immediate allocation."
        )
    top_team = max(teams, key=teams.get) if teams else "SEO"
    top_count = teams.get(top_team, 0)
    if top_count > 0:
        advisory_lines.append(
            f"Heaviest load: {top_team} team ({top_count} items) — "
            f"{'consider redistributing or extending the sprint' if top_count > 8 else 'manageable within sprint' if top_count > 4 else 'well-balanced'}."
        )
    if advisory_lines:
        _advisor_note(slide, MARGIN, y3 + 0.1, CONTENT_W, 0.55,
                      "Advisor's Analysis",
                      " | ".join(advisory_lines), PRIMARY_BLUE)

    add_footer(slide, page_num=page)
    return page


def _team_slide(prs: Presentation, team_name: str, items: list[ActionItem],
                accent: RGBColor, page: int, chunk_idx: int = 0,
                total_chunks: int = 1) -> int:
    """Team task table with advisor context."""
    page += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    subtitle = f"{len(items)} items" if chunk_idx == 0 else f"continued ({chunk_idx + 1}/{total_chunks})"
    add_header_bar(slide, team_name, subtitle, page)

    headers = ["#", "Task", "Pri", "Impact", "Effort", "Owner", "ETA"]
    col_widths = [0.3, 3.5, 0.4, 0.6, 0.5, 1.0, 0.8]
    total_w = sum(col_widths)
    n_cols = len(headers)
    n_rows = len(items) + 1
    row_h = 0.32

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(MARGIN), Inches(1.2),
        Inches(total_w), Inches(row_h * n_rows)
    )
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(hdr)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    has_risks = False
    for ri, item in enumerate(items):
        bg = GRAY_BG if ri % 2 == 0 else WHITE
        vals = [
            str(ri + 1),
            (item.task or "")[:60],
            item.priority or "P3",
            (item.impact or "medium")[:8],
            (item.effort or "—")[:8],
            (item.owner or "—")[:12],
            (item.eta or "—")[:10],
        ]
        for ci, val in enumerate(vals):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.02)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)

        if item.priority == "P1" and item.impact == "high":
            has_risks = True

    # Advisor note below table
    table_bottom = 1.2 + row_h * n_rows + 0.15
    if table_bottom < 5.5:
        notes = []
        p1_items = [a for a in items if a.priority == "P1"]
        p1_high_items = [a for a in items if a.priority == "P1" and a.impact == "high"]
        if p1_high_items:
            notes.append(
                f"{len(p1_high_items)} P1-high-impact item{'s' if len(p1_high_items) != 1 else ''} "
                f"— prioritize these first this sprint."
            )
        if p1_items and not p1_high_items:
            notes.append(
                f"{len(p1_items)} P1 item{'s' if len(p1_items) != 1 else ''} to address — "
                "no high-impact flags, but don't let them slip."
            )
        unowned = [a for a in items if not a.owner or a.owner == "Unassigned"]
        if unowned:
            notes.append(
                f"{len(unowned)} unassigned item{'s' if len(unowned) != 1 else ''} — "
                "assign owners to avoid blockers."
            )
        if notes:
            _advisor_note(slide, MARGIN, table_bottom, CONTENT_W, 0.45,
                          "Advisor's Note", " | ".join(notes[:2]), accent)

    add_footer(slide, page_num=page)
    return page


def _priority_matrix_slide(prs: Presentation, items: list[ActionItem],
                           page: int) -> int:
    """Advisor's priority-impact grid — what to do first, what to defer."""
    page += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_header_bar(slide, "Priority-Advisor Matrix",
                   "Impact vs urgency — advisor's recommendation on sequencing",
                   page, "Matrix")

    quadrants: dict[tuple[str, str], list[ActionItem]] = {}
    for item in items:
        prio = item.priority if item.priority in ("P1", "P2", "P3") else "P3"
        imp = item.impact.lower() if item.impact.lower() in ("high", "medium", "low") else "medium"
        quadrants.setdefault((prio, imp), []).append(item)

    y = 1.3
    for prio in ["P1", "P2", "P3"]:
        x = MARGIN
        pc = _priority_color(prio)
        _rect(slide, x, y, 0.8, 1.4, pc)
        _tb(slide, x + 0.05, y + 0.5, 0.7, 0.3,
            prio, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        x += 0.9
        for impact in ["high", "medium", "low"]:
            ic = _impact_color(impact)
            items_in_cell = quadrants.get((prio, impact), [])
            _rect(slide, x, y, 3.0, 1.4, GRAY_BG if items_in_cell else WHITE)
            _rect(slide, x, y, 3.0, 0.25, ic)
            _tb(slide, x + 0.05, y + 0.25, 2.9, 1.05,
                "\n".join(f"\u2022 {(item.task or '')[:35]}" for item in items_in_cell[:4])
                or "—",
                size=8, color=DARK_TEXT)
            x += 3.1
        y += 1.45

    # Advisor recommendation
    p1h = quadrants.get(("P1", "high"), [])
    p3l = quadrants.get(("P3", "low"), [])
    if p1h or p3l:
        rec_lines = []
        if p1h:
            rec_lines.append(
                f"Start with P1-high ({len(p1h)} items) — highest business impact and urgency."
            )
        if p3l:
            rec_lines.append(
                f"Defer P3-low ({len(p3l)} items) — lowest priority, revisit next sprint."
            )
        nh = 0.4 + 0.15 * (len(rec_lines) - 1)
        _advisor_note(slide, MARGIN, y + 0.05, CONTENT_W, nh,
                      "Advisor's Sequencing Recommendation",
                      " | ".join(rec_lines), PRIMARY_BLUE)

    add_footer(slide, page_num=page)
    return page


def _timeline_slide(prs: Presentation, items: list[ActionItem],
                    page: int) -> int:
    """Phased delivery timeline with advisor's phase guidance."""
    page += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_header_bar(slide, "Delivery Timeline",
                   "Phased approach — advisor's recommended sequence",
                   page, "Timeline")

    phases: list[tuple[str, str, list[ActionItem]]] = [
        ("Week 1", "Immediate", []),
        ("Week 2", "Short-term", []),
        ("Week 3-4", "This sprint", []),
        ("This month", "Medium-term", []),
        ("Next sprint", "Next cycle", []),
        ("Ongoing", "Continuous", []),
    ]
    phase_map = {p[0].lower(): (p[0], p[1], []) for p in phases}
    for item in items:
        eta = (item.eta or "").lower().strip()
        matched = False
        for key in phase_map:
            if key in eta:
                phase_map[key][2].append(item)
                matched = True
                break
        if not matched:
            phase_map["this month"][2].append(item)

    col_w = (CONTENT_W - 0.5) / 4
    y = 1.3
    col = 0
    for key, (phase_label, desc, phase_items) in sorted(phase_map.items(),
                                                         key=lambda x: list(phase_map.keys()).index(x[0])):
        x = MARGIN + col * (col_w + 0.15)
        if not phase_items:
            continue

        accent = [TEAL, PRIMARY_BLUE, AMBER, CORAL, GRAY_MID, NAVY][col % 6]
        _rect(slide, x, y, col_w, 0.35, accent)
        _tb(slide, x + 0.05, y + 0.02, col_w - 0.1, 0.16,
            phase_label, size=9, bold=True, color=WHITE)
        _tb(slide, x + 0.05, y + 0.18, col_w - 0.1, 0.15,
            desc, size=8, color=WHITE)

        y2 = y + 0.4
        for idx, item in enumerate(phase_items[:5]):
            pc = _priority_color(item.priority)
            _tb(slide, x + 0.05, y2, col_w - 0.1, 0.25,
                f"[{item.priority}] {(item.task or '')[:35]}",
                size=8, color=DARK_TEXT)
            y2 += 0.28

        col += 1
        if col >= 4:
            col = 0
            y = y2 + 0.15

    # Advisor phase tip
    first_phase = phase_map.get("week 1")
    if first_phase and first_phase[2]:
        p1_count = sum(1 for a in first_phase[2] if a.priority == "P1")
        tip = (
            f"Week 1 has {len(first_phase[2])} items ({p1_count} P1). "
            f"Ensure the team is staffed and blockers are cleared before sprint start."
        )
        _advisor_note(slide, MARGIN, y + 0.1, CONTENT_W, 0.4,
                      "Advisor's Phase Tip", tip, TEAL)

    add_footer(slide, page_num=page)
    return page


# ═══════════════════════════════════════════════════════════════
# MAIN BUILD FUNCTION
# ═══════════════════════════════════════════════════════════════

def build_action_plan_doc(facts: ReportFacts, agency: str = "",
                           client: str = "", month: str = "") -> bytes:
    """Build a 6-10 slide internal team action plan.

    Voice: Smart Advisor / Senior Analyst guiding the team.
    Slides: Cover, Task Landscape, Team breakdowns (with advisor notes),
    Priority-Advisor Matrix, Delivery Timeline.
    """
    teams: dict[str, list[ActionItem]] = {
        "SEO": [], "Dev": [], "Content": [], "Local": [],
    }
    for a in facts.action_plan:
        t = a.team if a.team in teams else "SEO"
        teams[t].append(a)

    for t in teams:
        teams[t].sort(key=lambda x: (PRIO_SORT.get(x.priority, 99), x.eta or "Z"))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    page = _cover_slide(prs, client, month, agency)

    all_items = [a for items in teams.values() for a in items]
    page = _overview_slide(prs, all_items, page)

    team_config = [
        ("SEO Team", teams["SEO"], TEAL),
        ("Dev Team", teams["Dev"], CORAL),
        ("Content Team", teams["Content"], AMBER),
        ("Local SEO", teams["Local"], NAVY),
    ]
    MAX_ROWS = 10
    for title, items, accent in team_config:
        if not items:
            continue
        chunks = [items[i:i + MAX_ROWS] for i in range(0, len(items), MAX_ROWS)]
        for chunk_idx, chunk in enumerate(chunks):
            page = _team_slide(prs, title, chunk, accent, page,
                               chunk_idx, len(chunks))

    page = _priority_matrix_slide(prs, all_items, page)
    page = _timeline_slide(prs, all_items, page)

    from report.ppt_design import validate_ppt_guardrail
    validate_ppt_guardrail(prs, label="Action Plan")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info(f"Action plan: {page} slides, {buf.getbuffer().nbytes // 1024} KB")
    return buf.getvalue()
