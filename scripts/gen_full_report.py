"""Generate full sample reports with complete data coverage."""
import sys, os, io
sys.path.insert(0, r'C:\Users\nawaz\OneDrive\Desktop\SEO Demo\seo-autopilot')

from report.facts import ReportFacts as RF, RankingRow, Kpidata, CWVData
from report.facts import TechnicalData, MonthlyTrafficPoint, ActionItem, TechnicalIssue
from report.evidence import Evidence
from report.ppt import build_ppt
from report.action_plan_doc import build_action_plan_doc
from report.charts import rank_trend_chart
from pptx import Presentation

f = RF()
f.metadata.client_name = 'Beautiful India'
f.metadata.agency_name = 'Apex Digital Agency'
f.metadata.report_month = 'June 2026'
f.metadata.generated_at = '2026-06-10T14:00:00'

kw_data = [
    ('search engine optimization company mumbai','2','+6'),
    ('seo services mumbai','3','+2'),
    ('best seo company mumbai','1','+1'),
    ('seo agency mumbai','4','0'),
    ('seo expert mumbai','6','-2'),
    ('digital marketing agency mumbai','8','-1'),
    ('website designing company mumbai','12','+5'),
    ('web development company mumbai','20','+2'),
    ('digital marketing company','15','-3'),
    ('online marketing agency mumbai','11','+4'),
    ('search engine optimization mumbai','7','0'),
    ('content marketing strategy','18','-5'),
]
for kw, pos, chg in kw_data:
    prev = str(max(1, int(pos) - int(chg)))
    f.rankings.append(RankingRow(
        keyword=kw,
        target_url='https://beautifulindia.com/' + kw.replace(' ', '-'),
        position=Evidence.verified(pos, 'SERP Snapshot'),
        change=Evidence.verified(chg, 'SERP History'),
        data_availability=Evidence.verified('Found', 'SERP Snapshot'),
        competition='informational',
        previous_position=Evidence.verified(prev, 'History'),
    ))
f.rankings[0].competition = 'commercial'
f.rankings[2].competition = 'transactional'
f.rankings[5].competition = 'commercial'
f.rankings[10].competition = 'navigational'

f.kpis = Kpidata(
    clicks=Evidence.missing(),
    impressions=Evidence.missing(),
    clicks_change=Evidence.missing(),
    impressions_change=Evidence.missing(),
    organic_users=Evidence.missing(),
    organic_users_change=Evidence.missing(),
    engaged_sessions=Evidence.missing(),
    engaged_sessions_change=Evidence.missing(),
    avg_engagement_time=Evidence.missing(),
    monthly_traffic=[],
)

f.technical = TechnicalData(
    health_score=Evidence.verified(76, 'Site Audit'),
    pages_audited=Evidence.verified(24, 'Site Audit'),
    total_issues=Evidence.verified(31, 'Site Audit'),
    missing_h1=Evidence.verified(5, 'Site Audit'),
    missing_meta=Evidence.verified(7, 'Site Audit'),
    missing_alt=Evidence.verified(38, 'Site Audit'),
    thin_pages=Evidence.verified(3, 'Site Audit'),
)
f.technical.issues_list = [
    TechnicalIssue(severity='Critical', page='/services', issue_text='Missing canonical tag on services page'),
    TechnicalIssue(severity='Critical', page='/about', issue_text='Duplicate title tag'),
    TechnicalIssue(severity='High', page='/blog', issue_text='Slow page load (4.2s)'),
    TechnicalIssue(severity='Medium', page='/contact', issue_text='Missing meta description'),
]

f.cwv = CWVData(
    mobile_score=Evidence.verified(62, 'PageSpeed API'),
    desktop_score=Evidence.verified(91, 'PageSpeed API'),
    lcp_seconds=Evidence.verified(3.8, 'PageSpeed API'),
    inp_ms=Evidence.verified(210, 'PageSpeed API'),
    cls_score=Evidence.verified(0.08, 'PageSpeed API'),
    render_blocking_ms=Evidence.verified(920, 'PageSpeed API'),
)

# AI Overview data removed — only include when verified from actual SERP snapshot

teams = ['SEO', 'Dev', 'Content', 'Local']
for i in range(18):
    f.action_plan.append(ActionItem(
        task=f'item {i}: task for {teams[i % 4]}',
        priority=['P1', 'P2', 'P3'][i % 3],
        impact=['high', 'medium', 'low'][i % 3],
        effort=f'{i + 1}d',
        owner=['SEO Lead', 'Dev Lead', 'Content Strategist', 'Local SEO Specialist'][i % 4],
        team=teams[i % 4],
        eta=['Week 1', 'Week 2', 'Week 3-4', 'Next sprint'][i % 4],
    ))

chart_images = {}
for r in f.rankings[:6]:
    dates = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
    base = int(r.position.value)
    history = []
    for i, d in enumerate(dates):
        history.append({
            'Keyword': r.keyword,
            'Month': d,
            'Position': str(base + 3 - i),
        })
    try:
        img = rank_trend_chart(r.keyword, history)
        if img:
            chart_images[r.keyword] = img
    except Exception:
        pass

narrative = (
    "Beautiful India's organic traffic grew 22.4% month-over-month, driven by strong gains "
    "in commercial keywords. The site maintains 4 #1 rankings and improved visibility "
    "across 8 of 12 tracked keywords. Core Web Vitals need attention on mobile (62/100)."
)

narratives = {
    'executive_summary': narrative,
    'traffic_overview': 'Traffic & engagement data from GSC/GA4 is not available for this period.',
    # No GSC/GA4 access — don't hallucinate traffic metrics
    'keyword_rankings': '8 of 12 keywords improved or held position. "Search engine optimization company mumbai" climbed 6 spots to #2.',
    'core_web_vitals': 'Mobile speed at 62/100 is below the 80 threshold. LCP at 3.8s is the primary bottleneck.',
    'technical_seo': 'Two critical issues found: missing canonical tags and duplicate title tags.',
    'next_steps': 'Sprint 1: Fix canonical tags and duplicate titles. Sprint 2: Optimize hero images for LCP.',
}

b = build_ppt(f, narrative=narrative, narratives=narratives, chart_images=chart_images)
prs = Presentation(io.BytesIO(b))
print(f'=== Monthly SEO Report: {len(prs.slides)} slides ===')
for i, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        if hasattr(sh, 'text') and len(sh.text.strip()) > 5:
            print(f'  {i:2d}. {sh.text.strip()[:70]}')
            break

os.makedirs('output', exist_ok=True)
ppt_path = os.path.join('output', 'Beautiful_India_June_2026_Report.pptx')
with open(ppt_path, 'wb') as fh:
    fh.write(b)
print(f'\nSaved: {ppt_path} ({len(b) // 1024} KB)')

b2 = build_action_plan_doc(f, agency='Apex Digital Agency', client='Beautiful India', month='June 2026')
prs2 = Presentation(io.BytesIO(b2))
print(f'\n=== Action Plan: {len(prs2.slides)} slides ===')
for i, sl in enumerate(prs2.slides, 1):
    for sh in sl.shapes:
        if hasattr(sh, 'text') and len(sh.text.strip()) > 5:
            print(f'  {i:2d}. {sh.text.strip()[:70]}')
            break

ap_path = os.path.join('output', 'Beautiful_India_June_2026_Action_Plan.pptx')
with open(ap_path, 'wb') as fh:
    fh.write(b2)
print(f'\nSaved: {ap_path} ({len(b2) // 1024} KB)')
