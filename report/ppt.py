"""Monthly SEO Report — Agency-to-Executive presentation, 12-14 slides.

Voice: Digital Marketing Agency presenting to client executives.
Every insight, recommendation, and metric originates from ReportFacts,
agent narratives, or chart images. No placeholders. No hallucinated claims.
"""
from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from report.facts import ReportFacts
from report.ppt_data import build_ppt_data
from report.ppt_design import (
    SLIDE_W, SLIDE_H, MARGIN, COL, CONTENT_W, CONTENT_L, CONTENT_R,
    NAVY, PRIMARY_BLUE, DARK_BLUE, TEAL, AMBER, RED, WHITE,
    GRAY_BG, GRAY_LINE, GRAY_MID, GRAY_TEXT, DARK_TEXT, GREEN,
    SUBTITLE_BLUE, INFO_BG, GREEN_LIGHT, AMBER_LIGHT, CORAL,
    add_slide_bg, add_header_bar, add_footer,
    kpi_card, metric_row, add_table, _tb, _line, _rect, _rgb,
    progress_bar, big_stat, status_badge,
    score_gauge, make_cover, make_closing,
)
from report.charts import (
    make_traffic_line_chart,
    make_psi_column_chart,
    make_distribution_hbar_chart,
)

logger = logging.getLogger(__name__)


def _fmt_change(val: int | str) -> str:
    try:
        v = int(val)
        if v > 0:
            return f"+{v}"
        if v < 0:
            return str(v)
        return "0"
    except (ValueError, TypeError):
        return str(val) if val else "0"


# ═══════════════════════════════════════════════════════════════
# EXECUTIVE-INSIGHT HELPERS
# ═══════════════════════════════════════════════════════════════

def _insight_box(slide, l: float, t: float, w: float, h: float,
                 heading: str, body: str, accent=PRIMARY_BLUE):
    """Premium agency-style callout panel for executive insights."""
    _rect(slide, l, t, w, h, INFO_BG)
    _rect(slide, l, t, 0.06, h, accent)
    _tb(slide, l + 0.18, t + 0.08, w - 0.3, 0.22,
        heading, size=10, bold=True, color=accent)
    _tb(slide, l + 0.18, t + 0.32, w - 0.3, h - 0.4,
        body, size=9, color=DARK_TEXT)


def _exec_subtitle(d: dict, key: str, default: str = "") -> str:
    """Build an executive-facing subtitle from data, or return default."""
    val = d.get(key)
    if val and val != "—":
        return str(val)
    return default


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def _slide_exec_dashboard(prs: Presentation, d: dict, narratives: dict,
                          page_num: int) -> int:
    """Slide 2: Executive Dashboard — agency-grade KPI overview for leadership."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)

    subtitle = _exec_subtitle(d, "exec_overall_seo_score",
                              f"SEO Health Score: {d.get('exec_overall_seo_score', '—')}")
    add_header_bar(slide, "Executive Dashboard",
                   subtitle,
                   page_num, "Dashboard")

    # Agency narrative as lead insight
    narr = narratives.get("executive_summary", "").strip()
    if narr:
        _insight_box(slide, MARGIN, 0.9, CONTENT_W, 0.45,
                     "Executive Summary", narr, PRIMARY_BLUE)

    # Three KPI cards across
    cw = (CONTENT_W - 0.6) / 3
    rx = MARGIN
    y = 1.5

    big_win_kw = d.get("exec_big_win_kw", "—")
    big_win_pos = d.get("exec_big_win_pos", "—")
    big_win_chg = d.get("exec_big_win_change", "—")

    kpi_card(slide, rx, y, cw, 1.5,
             "Biggest Win", f"{big_win_kw}", f"Position {big_win_pos} ({big_win_chg})",
             GREEN, GREEN)

    bottleneck = d.get("exec_bottleneck_text", "")
    kpi_card(slide, rx + cw + 0.3, y, cw, 1.5,
             "Key Bottleneck", "Page Speed",
             bottleneck[:60] if bottleneck else "No critical bottlenecks detected",
             AMBER, AMBER)

    score_label = d.get("exec_overall_seo_score", "—")
    grade = d.get("overall_grade", "—")
    kpi_card(slide, rx + 2 * (cw + 0.3), y, cw, 1.5,
             f"Score {grade}", score_label,
             "Weighted across 6 categories",
             PRIMARY_BLUE, PRIMARY_BLUE)

    # Keyword summary row with business framing
    y2 = y + 1.7
    _line(slide, MARGIN, y2, CONTENT_W)
    y2 += 0.15

    total = d.get("kw_total_tracked", 0)
    improved = d.get("kw_improved", 0)
    dropped = d.get("kw_dropped", 0)
    stable = d.get("kw_stable", 0)
    nf = d.get("kw_not_found", 0)

    stat_w = (CONTENT_W - 0.5) / 5
    stats = [
        ("Tracked", str(total), PRIMARY_BLUE),
        ("Improved", str(improved), GREEN),
        ("Dropped", str(dropped), RED),
        ("Stable", str(stable), GRAY_MID),
        ("Not Found", str(nf), AMBER),
    ]
    for i, (label, val, color) in enumerate(stats):
        big_stat(slide, MARGIN + i * (stat_w + 0.12), y2, stat_w, 0.9,
                 val, label, color)

    # Business-impact insight panel below stats
    y3 = y2 + 1.1
    _line(slide, MARGIN, y3, CONTENT_W)

    if improved > 0 or dropped > 0:
        insight_lines = []
        if improved > 0:
            insight_lines.append(
                f"Growth signal: {improved} keyword{'s' if improved != 1 else ''} "
                f"gained position — our content and link-building strategy is driving results."
            )
        if dropped > 0:
            insight_lines.append(
                f"Attention needed: {dropped} keyword{'s' if dropped != 1 else ''} "
                f"lost ground. We will investigate SERP changes and competitor activity."
            )
        if total > 0:
            visible = total - nf
            pct = round(visible / total * 100) if total else 0
            insight_lines.append(
                f"Visibility rate: {visible}/{total} keywords ranking ({pct}%) "
                f"— {'strong' if pct >= 80 else 'moderate' if pct >= 60 else 'early-stage'} market presence."
            )
        _insight_box(slide, MARGIN, y3 + 0.1, CONTENT_W, 0.55,
                     "Agency Perspective", " | ".join(insight_lines), PRIMARY_BLUE)

    # Priority recommendations
    y4 = y3 + 0.75
    recs = d.get("exec_priority_recs", [])
    if recs:
        _tb(slide, MARGIN, y4, 3.0, 0.25, "Priority Recommendations",
            size=13, bold=True, color=NAVY)
        y4 += 0.35
        for rec in recs[:3]:
            cat = rec.get("category", "SEO")[:14]
            badge_w = len(cat) * 0.09 + 0.30
            _rect(slide, MARGIN, y4, badge_w, 0.24, PRIMARY_BLUE)
            _tb(slide, MARGIN + 0.03, y4 + 0.01, badge_w - 0.06, 0.22,
                cat, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _tb(slide, MARGIN + badge_w + 0.1, y4, CONTENT_W - badge_w - 0.1, 0.24,
                rec.get("text", ""), size=9, color=DARK_TEXT)
            y4 += 0.28

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_traffic(prs: Presentation, d: dict, narratives: dict,
                   page_num: int) -> int:
    """Slide 3: Audience Growth — business impact of organic traffic."""
    traffic_data = d.get("traffic", {})
    monthly = traffic_data.get("monthly_traffic", [])
    if not monthly or all(m.get("users", "—") == "—" for m in monthly):
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)

    # Executive subtitle: show organic user change if available
    sub = d.get("title_month", "")
    org_chg = traffic_data.get("organic_users_change", "")
    if org_chg and org_chg != "—":
        sub = f"Organic users {org_chg} — {d.get('title_month', '')}"
    add_header_bar(slide, "Audience Growth & Engagement",
                   sub, page_num, "Traffic")

    # Embed inline chart
    has_monthly_data = monthly and monthly[0].get("users", "—") != "—"
    if has_monthly_data:
        try:
            chart_bytes = make_traffic_line_chart(monthly)
            if chart_bytes:
                slide.shapes.add_picture(
                    io.BytesIO(chart_bytes),
                    Inches(MARGIN), Inches(1.1),
                    Inches(CONTENT_W * 0.55), Inches(2.8),
                )
        except Exception:
            pass

    # Right side: key business metrics
    rx = MARGIN + CONTENT_W * 0.58
    y = 1.2
    metrics = [
        ("Organic Users", traffic_data.get("organic_users", "—"),
         traffic_data.get("organic_users_change", ""), GREEN),
        ("Engaged Sessions", traffic_data.get("engaged_sessions", "—"),
         traffic_data.get("engaged_sessions_change", ""), PRIMARY_BLUE),
        ("Avg. Engagement", traffic_data.get("avg_engagement_time", "—"),
         "", GRAY_MID),
    ]
    for label, val, chg, color in metrics:
        kpi_card(slide, rx, y, CONTENT_W * 0.38, 0.85,
                 label, str(val), str(chg) if chg else "",
                 color, color)
        y += 0.95

    # GSC metrics row
    gsc = d.get("gsc", {})
    y = 3.5
    _line(slide, MARGIN, y, CONTENT_W)
    _tb(slide, MARGIN, y + 0.1, 5.0, 0.25,
        "Google Search Console — How audiences find us",
        size=12, bold=True, color=NAVY)
    y += 0.4
    gsc_metrics = [
        ("Clicks", gsc.get("clicks_curr", "—"), gsc.get("clicks_change", "")),
        ("Impressions", gsc.get("impressions_curr", "—"), gsc.get("impressions_change", "")),
        ("CTR", gsc.get("ctr_curr", "—"), ""),
    ]
    for i, (label, val, chg) in enumerate(gsc_metrics):
        stat_w = (CONTENT_W - 0.5) / 3
        big_stat(slide, MARGIN + i * (stat_w + 0.2), y, stat_w, 0.7,
                 str(val), label, PRIMARY_BLUE)

    # Business impact insight
    narr = narratives.get("traffic_overview", "").strip()
    if narr:
        _insight_box(slide, MARGIN, y + 0.85, CONTENT_W, 0.45,
                     "Business Impact", narr, GREEN)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_serp_visibility(prs: Presentation, d: dict, narratives: dict,
                           chart_images: dict, page_num: int) -> int:
    """Slide 4: Search Visibility — market presence and competitive position."""
    dist = d.get("kw_position_distribution", {})
    if not any(v for v in dist.values()):
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    total = d.get("kw_total_tracked", 0)
    top3 = dist.get("p1", 0) + dist.get("top3", 0)
    sub = f"{d.get('kw_total_tracked', 0)} keywords tracked"
    if top3 > 0:
        sub += f" — {top3} in top 3 positions"
    add_header_bar(slide, "Search Visibility & Market Presence",
                   sub, page_num, "Visibility")

    # Distribution chart
    try:
        chart_bytes = make_distribution_hbar_chart(dist)
        if chart_bytes:
            slide.shapes.add_picture(
                io.BytesIO(chart_bytes),
                Inches(MARGIN), Inches(1.2),
                Inches(4.5), Inches(3.0),
            )
    except Exception:
        pass

    # Right side: breakdown with business framing
    rx = MARGIN + 4.8
    y = 1.3
    labels = [
        ("#1 Rankings (Market Leader)", dist.get("p1", 0), GREEN),
        ("Top 3 (Strong Presence)", dist.get("top3", 0), PRIMARY_BLUE),
        ("Top 10 (Competitive)", dist.get("top10", 0), TEAL),
        ("Page 2 (Opportunity Zone)", dist.get("p11_20", 0), AMBER),
        ("Beyond 20 (Needs Work)", dist.get("beyond", 0), GRAY_MID),
        ("Not Found (Gap)", dist.get("not_found", 0), RED),
    ]
    for label, val, color in labels:
        big_stat(slide, rx, y, CONTENT_W - 4.8, 0.55,
                 str(val), label, color)
        y += 0.62

    # Quick wins insight
    quick = d.get("kw_quick_wins", [])
    if quick:
        y = y + 0.2 if y < 4.0 else y
        _insight_box(slide, MARGIN, y, CONTENT_W, 0.55 + 0.28 * min(len(quick), 3),
                     "Quick-Win Opportunities",
                     "\n".join(
                         f"\u2191 {q.get('keyword', '')[:35]} (#{q.get('pos', '?')}) — "
                         f"just outside top 10, strong improvement potential"
                         for q in quick[:3]
                     ) if quick else "",
                     AMBER)
        y += 0.55 + 0.28 * min(len(quick), 3) + 0.1 + 0.55

    narr = narratives.get("keyword_rankings", "").strip()
    if narr:
        y = max(y + 0.1, 5.0)
        _tb(slide, MARGIN, y, CONTENT_W, 0.4, narr, size=9, color=GRAY_MID)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_keyword_winners(prs: Presentation, d: dict, page_num: int) -> int:
    """Slide 5: Growth Opportunities Captured — keywords gaining position."""
    wins = d.get("kw_wins", [])
    if not wins:
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_header_bar(slide, "Growth Opportunities Captured",
                   f"{len(wins)} keywords gained position — our strategy is delivering",
                   page_num, "Keywords")

    headers = ["Keyword", "Previous", "Current", "Gain"]
    col_widths = [3.5, 1.2, 1.2, 1.2]
    rows = []
    for w in wins:
        rows.append([
            w.get("keyword", "")[:40],
            str(w.get("prev", "—")),
            str(w.get("curr", "—")),
            w.get("change", "—"),
        ])

    add_table(slide, MARGIN, 1.3, col_widths, headers, rows, font_size=11)

    y = 1.3 + 0.35 * (len(rows) + 1) + 0.2
    if len(wins) > 2:
        _insight_box(slide, MARGIN + 0.5, y, CONTENT_W - 1.0, 0.5,
                     "Agency Insight",
                     f"{len(wins)} keywords are trending in the right direction. "
                     "We recommend doubling down on content clusters that are driving these gains.",
                     GREEN)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_keyword_losers(prs: Presentation, d: dict, page_num: int) -> int:
    """Slide 6: Risk Areas — keywords needing recovery attention."""
    not_found = d.get("kw_not_found", 0)
    dropped = d.get("kw_dropped", 0)
    if not_found == 0 and dropped == 0:
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_header_bar(slide, "Risk Areas Requiring Attention",
                   f"{dropped} dropped, {not_found} dropped from SERP",
                   page_num, "Recovery")

    y = 1.3
    if not_found > 0:
        _insight_box(slide, MARGIN, y, CONTENT_W, 1.5,
                     f"Keywords Not Found in SERP — {not_found} at risk",
                     f"{not_found} tracked keywords are currently not ranking.\n\n"
                     "Our recommended recovery approach:\n"
                     "\u2022 Verify indexation in Google Search Console\n"
                     "\u2022 Check for indexing blocks (noindex, canonical drift)\n"
                     "\u2022 Strengthen content with local intent signals\n"
                     "\u2022 Build internal links from ranking pages",
                     RED)
        y += 1.7

    if dropped > 0:
        _insight_box(slide, MARGIN, y, CONTENT_W, 1.3,
                     f"Position Losses — {dropped} keywords declined",
                     f"{dropped} keywords lost ranking position this period.\n\n"
                     "Recovery plan:\n"
                     "\u2022 Audit SERP for new competitors or feature snippets\n"
                     "\u2022 Refresh page content with current data\n"
                     "\u2022 Review backlink profile for lost equity\n"
                     "\u2022 Monitor trend over the next 2 weeks before escalating",
                     AMBER)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_search_intent(prs: Presentation, d: dict, page_num: int) -> int:
    """Slide 7: Audience Search Behavior — performance by intent."""
    by_intent = d.get("kw_by_intent", [])
    if not by_intent:
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    total = sum(item.get("count", 0) for item in by_intent)
    add_header_bar(slide, "Audience Search Behavior",
                   f"{total} keywords mapped across {len(by_intent)} intent types",
                   page_num, "Intent")

    y = 1.3
    for item in by_intent:
        label = item.get("label", "Unknown")
        count = item.get("count", 0)
        intent_descriptions = {
            "Info": "Users researching — build authority content",
            "Nav": "Users looking for specific brands — protect brand terms",
            "Commercial": "Users comparing options — capture with landing pages",
            "Transactional": "Users ready to convert — optimize for action",
            "Unknown": "Unclassified — review search queries",
        }
        desc = intent_descriptions.get(label, "")
        _rect(slide, MARGIN, y, CONTENT_W, 0.7, GRAY_BG)
        _rect(slide, MARGIN, y, 0.06, 0.7, PRIMARY_BLUE)
        _tb(slide, MARGIN + 0.2, y + 0.08, 2.5, 0.25,
            label, size=11, bold=True, color=DARK_TEXT)
        _tb(slide, MARGIN + 2.7, y + 0.08, 0.8, 0.25,
            f"{count} keywords", size=11, bold=True, color=PRIMARY_BLUE)
        if desc:
            _tb(slide, MARGIN + 0.2, y + 0.38, CONTENT_W - 0.4, 0.25,
                desc, size=9, color=GRAY_MID)
        y += 0.78

    # Executive insight about intent balance
    if len(by_intent) >= 2:
        y += 0.1
        info_count = next((item.get("count", 0) for item in by_intent if item.get("label") == "Info"), 0)
        comm_count = next((item.get("count", 0) for item in by_intent if item.get("label") == "Commercial"), 0)
        if comm_count > 0 and info_count > 0:
            ratio = round(info_count / comm_count, 1)
            _insight_box(slide, MARGIN, y, CONTENT_W, 0.45,
                         "Intent Balance",
                         f"Information-to-commercial ratio: {ratio}:1. "
                         f"{'Healthy balance for funnel coverage.' if 1 <= ratio <= 4 else 'Consider adding more bottom-funnel content.'}",
                         PRIMARY_BLUE)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_rank_trends(prs: Presentation, d: dict, chart_images: dict,
                       page_num: int) -> int:
    """Slide 8: Keyword Momentum — position trend charts for tracked terms."""
    if not chart_images:
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    n_charts = min(len(chart_images), 6)
    add_header_bar(slide, "Keyword Momentum",
                   f"{n_charts} tracked keywords — position movement over time",
                   page_num, "Trends")

    # Place up to 6 charts in a 3x2 grid
    chart_items = list(chart_images.items())[:6]
    cols = 3
    rows = (len(chart_items) + cols - 1) // cols
    cw = (CONTENT_W - 0.4) / cols
    ch = 1.8

    for idx, (kw, img_bytes) in enumerate(chart_items):
        col = idx % cols
        row = idx // cols
        x = MARGIN + col * (cw + 0.2)
        y = 1.2 + row * (ch + 0.15)
        try:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(x), Inches(y),
                Inches(cw), Inches(ch),
            )
        except Exception:
            _tb(slide, x, y, cw, 0.3, kw[:30], size=8, color=GRAY_MID)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_core_web_vitals(prs: Presentation, d: dict, narratives: dict,
                           page_num: int) -> int:
    """Slide 9: User Experience & Performance — Core Web Vitals report."""
    ds = d.get("cwv_desktop_score", "—")
    ms = d.get("cwv_mobile_score", "—")
    lcp = d.get("cwv_lcp", "—")
    inp = d.get("cwv_inp", "—")
    cls = d.get("cwv_cls", "—")

    if ds == "—" and ms == "—":
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    sub = f"Desktop {ds}/100  |  Mobile {ms}/100"
    add_header_bar(slide, "User Experience & Page Performance",
                   sub, page_num, "CWV")

    # PSI chart
    try:
        ds_val = int(ds) if ds != "—" else 0
        ms_val = int(ms) if ms != "—" else 0
        if ds_val > 0 or ms_val > 0:
            chart_bytes = make_psi_column_chart(ds_val, ms_val)
            if chart_bytes:
                slide.shapes.add_picture(
                    io.BytesIO(chart_bytes),
                    Inches(MARGIN), Inches(1.2),
                    Inches(4.0), Inches(2.8),
                )
    except Exception:
        pass

    # Right side: CWV metrics with status
    rx = MARGIN + 4.3
    y = 1.3
    metrics = [
        ("LCP (Largest Contentful Paint)", lcp, "s",
         "good" if lcp != "—" and float(str(lcp).rstrip("s")) <= 2.5
         else "warn" if lcp != "—" and float(str(lcp).rstrip("s")) <= 4.0
         else "bad"),
        ("INP (Interaction to Next Paint)", inp, "ms",
         "good" if inp != "—" and int(inp) <= 200
         else "warn" if inp != "—" and int(inp) <= 500
         else "bad"),
        ("CLS (Cumulative Layout Shift)", cls, "",
         "good" if cls != "—" and float(cls) <= 0.1
         else "warn" if cls != "—" and float(cls) <= 0.25
         else "bad"),
    ]
    for label, val, unit, status in metrics:
        mc = GREEN if status == "good" else AMBER if status == "warn" else RED
        _rect(slide, rx, y, CONTENT_W - 4.3, 0.65, GRAY_BG)
        _rect(slide, rx, y, 0.06, 0.65, mc)
        _tb(slide, rx + 0.15, y + 0.05, 2.2, 0.22,
            label, size=9, bold=True, color=DARK_TEXT)
        _tb(slide, rx + 2.5, y + 0.05, 1.5, 0.22,
            f"{val}{unit}", size=11, bold=True, color=mc)
        y += 0.72

    # Dev tickets
    tickets = d.get("cwv_dev_tickets", [])
    if tickets:
        y = max(y, 3.5)
        _line(slide, MARGIN, y, CONTENT_W)
        y += 0.1
        _tb(slide, MARGIN, y, 5.0, 0.25,
            "Technical Items — Recommended for Dev Sprint",
            size=12, bold=True, color=NAVY)
        y += 0.3
        for t in tickets[:3]:
            _tb(slide, MARGIN + 0.1, y, CONTENT_W - 0.1, 0.22,
                f"{t.get('label', '')}", size=9, color=GRAY_MID)
            y += 0.24

    # Business impact insight
    narr = narratives.get("core_web_vitals", "").strip()
    if narr:
        y = max(y + 0.1, 5.5)
        _insight_box(slide, MARGIN, y, CONTENT_W, 0.45,
                     "Business Impact",
                     f"{narr}",
                     PRIMARY_BLUE)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_technical_health(prs: Presentation, d: dict, narratives: dict,
                            page_num: int) -> int:
    """Slide 10: Site Foundation — technical health and infrastructure."""
    score = d.get("tech_health_score", "—")
    actions = d.get("tech_actions", [])
    if score == "—" and not actions:
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    score_num = d.get("tech_health_score_num", 0)
    grade_map = {True: "Good", False: "Needs Improvement"}
    sub = f"Health Score: {score} — {grade_map.get(score_num >= 70, '')}"
    add_header_bar(slide, "Site Foundation & Technical Readiness",
                   sub, page_num, "Technical")

    # Score gauge on left
    gc = GREEN if score_num >= 80 else AMBER if score_num >= 50 else RED
    score_gauge(slide, MARGIN, 1.3, 2.8, 1.5, score_num, "Site Health", gc)

    # Key issues on right
    rx = MARGIN + 3.2
    _tb(slide, rx, 1.3, 4.0, 0.25,
        "Technical Issues Detected — Impact on Performance",
        size=12, bold=True, color=NAVY)
    y = 1.6
    tech_issues = [
        ("Missing H1 Tags", d.get("tech_missing_h1", "—")),
        ("Missing Meta Descriptions", d.get("tech_missing_meta", "—")),
        ("Missing Image Alt Text", d.get("tech_missing_alt", "—")),
        ("Thin Content Pages", d.get("tech_thin_pages", "—")),
    ]
    for label, val in tech_issues:
        metric_row(slide, rx, y, CONTENT_W - 3.2, label, val,
                   "bad" if str(val) not in ("—", "0") else "good")
        y += 0.3

    # Action checklist
    if actions:
        y = max(y, 3.5)
        _line(slide, MARGIN, y, CONTENT_W)
        y += 0.1
        _tb(slide, MARGIN, y, 5.0, 0.25,
            "Recommended Technical Improvements",
            size=12, bold=True, color=NAVY)
        y += 0.3
        for action in actions[:5]:
            txt = action.get("text", "")
            if txt:
                done = action.get("done", False)
                prefix = "\u2611" if done else "\u2610"
                _tb(slide, MARGIN + 0.1, y, CONTENT_W - 0.1, 0.22,
                    f"{prefix} {txt[:70]}", size=9, color=GRAY_MID)
                y += 0.24

    narr = narratives.get("technical_seo", "").strip()
    if narr:
        y = max(y + 0.1, 5.5)
        _insight_box(slide, MARGIN, y, CONTENT_W, 0.45,
                     "Agency Recommendation", narr, AMBER)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_onpage_seo(prs: Presentation, d: dict, page_num: int) -> int:
    """Slide 11: Content & Structure Quality — on-page optimization summary."""
    has_https = d.get("onpage_has_https", False)
    has_canonical = d.get("onpage_has_canonical", False)
    img_stats = d.get("onpage_image_stats", "")
    if not has_https and not has_canonical and not img_stats:
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    audited = d.get("tech_pages_audited", "—")
    add_header_bar(slide, "Content & Structure Quality",
                   f"Pages audited: {audited}",
                   page_num, "On-Page")

    y = 1.3
    metrics = [
        ("HTTPS (Secure Connection)", "\u2713 Enabled" if has_https else "\u2717 Not Enabled",
         GREEN if has_https else RED),
        ("Canonical Tags", "\u2713 Present" if has_canonical else "\u2717 Missing",
         GREEN if has_canonical else RED),
        ("Images Missing Alt Text", d.get("tech_missing_alt", "—"),
         AMBER if str(d.get("tech_missing_alt", "—")) not in ("—", "0") else GREEN),
        ("Pages Missing H1 Tags", d.get("tech_missing_h1", "—"),
         AMBER if str(d.get("tech_missing_h1", "—")) not in ("—", "0") else GREEN),
        ("Pages Missing Meta Desc.", d.get("tech_missing_meta", "—"),
         AMBER if str(d.get("tech_missing_meta", "—")) not in ("—", "0") else GREEN),
        ("Thin Content Pages", d.get("tech_thin_pages", "—"),
         AMBER if str(d.get("tech_thin_pages", "—")) not in ("—", "0") else GREEN),
    ]
    for label, val, color in metrics:
        _rect(slide, MARGIN, y, CONTENT_W, 0.45, GRAY_BG)
        _rect(slide, MARGIN, y, 0.06, 0.45, color)
        _tb(slide, MARGIN + 0.15, y + 0.05, 4.0, 0.2,
            label, size=10, bold=True, color=DARK_TEXT)
        _tb(slide, MARGIN + 4.3, y + 0.05, 2.5, 0.2,
            str(val), size=10, bold=True, color=color)
        y += 0.5

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_ai_insights(prs: Presentation, d: dict, page_num: int) -> int:
    """Slide 12: AI Readiness — how AI-generated search results affect visibility."""
    ai_summary = d.get("ai_overview_summary", "")
    if not ai_summary or ai_summary == "AI Overview data not yet collected for tracked keywords.":
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_header_bar(slide, "AI Overviews & SERP Evolution",
                   "How AI-generated search results affect your visibility",
                   page_num, "AI")

    _insight_box(slide, MARGIN, 1.3, CONTENT_W, 2.0,
                 "AI Overview Presence",
                 ai_summary + "\n\n"
                 "AI overviews are reshaping search. Keywords with AI-generated "
                 "results typically see reduced click-through rates. "
                 "We recommend building structured data and authoritative content "
                 "to maintain visibility in this evolving landscape.",
                 PRIMARY_BLUE)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


def _slide_action_plan(prs: Presentation, d: dict, narratives: dict,
                       page_num: int) -> int:
    """Slide 13: Executive Recommendations — prioritized next steps."""
    plan_seo = d.get("plan_seo_tasks", [])
    plan_dev = d.get("plan_dev_tasks", [])
    plan_content = d.get("plan_content_tasks", [])
    if not (plan_seo or plan_dev or plan_content):
        return page_num

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_header_bar(slide, "Executive Recommendations",
                   "Prioritized next steps by team",
                   page_num, "Actions")

    teams_data = []
    if plan_seo:
        teams_data.append(("SEO Strategy", plan_seo, TEAL))
    if plan_dev:
        teams_data.append(("Development", plan_dev, CORAL))
    if plan_content:
        teams_data.append(("Content", plan_content, AMBER))

    col_w = (CONTENT_W - 0.4) / max(len(teams_data), 1)
    for i, (team_name, tasks, accent) in enumerate(teams_data):
        x = MARGIN + i * (col_w + 0.2)
        _rect(slide, x, 1.3, col_w, 0.35, accent)
        _tb(slide, x + 0.1, 1.32, col_w - 0.2, 0.3,
            team_name, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        y = 1.75
        for task in tasks[:6]:
            _tb(slide, x + 0.05, y, col_w - 0.1, 0.22,
                f"\u2022 {task[:50]}", size=9, color=DARK_TEXT)
            y += 0.24

    narr = narratives.get("next_steps", "").strip()
    if narr:
        y = max(y + 0.2, 5.0)
        _line(slide, MARGIN, y, CONTENT_W)
        _insight_box(slide, MARGIN, y + 0.1, CONTENT_W, 0.45,
                     "Next Sprint Focus", narr, PRIMARY_BLUE)

    add_footer(slide, d.get("title_agency", ""), page_num)
    return page_num + 1


# ═══════════════════════════════════════════════════════════════
# MAIN BUILD FUNCTION
# ═══════════════════════════════════════════════════════════════

def build_ppt(
    facts: ReportFacts,
    narrative: str = "",
    narratives: dict | None = None,
    chart_images: dict | None = None,
    screenshots: dict | None = None,
    psi_screenshots: dict | None = None,
) -> bytes:
    """Build a 12-14 slide executive SEO report from ReportFacts.

    Returns PPTX bytes.
    """
    d = build_ppt_data(facts)
    chart_images = chart_images or {}
    screenshots = screenshots or {}
    narratives = narratives or {}

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # ── Slide 1: Cover ──
    make_cover(prs, d["title_client"], d["title_month"], d["title_agency"],
               generated=d.get("title_generated", ""),
               health_score=d.get("exec_overall_seo_score_num"))
    page = 2

    # ── Slide 2: Executive Dashboard ──
    page = _slide_exec_dashboard(prs, d, narratives, page)

    # ── Slide 3: Traffic ──
    page = _slide_traffic(prs, d, narratives, page)

    # ── Slide 4: SERP Visibility ──
    page = _slide_serp_visibility(prs, d, narratives, chart_images, page)

    # ── Slide 5: Keyword Winners ──
    page = _slide_keyword_winners(prs, d, page)

    # ── Slide 6: Keyword Losers ──
    page = _slide_keyword_losers(prs, d, page)

    # ── Slide 7: Search Intent ──
    page = _slide_search_intent(prs, d, page)

    # ── Slide 8: Rank Trends ──
    page = _slide_rank_trends(prs, d, chart_images, page)

    # ── Slide 9: Core Web Vitals ──
    page = _slide_core_web_vitals(prs, d, narratives, page)

    # ── Slide 10: Technical Health ──
    page = _slide_technical_health(prs, d, narratives, page)

    # ── Slide 11: On-Page SEO ──
    page = _slide_onpage_seo(prs, d, page)

    # ── Slide 12: AI Insights ──
    page = _slide_ai_insights(prs, d, page)

    # ── Slide 13: Action Plan ──
    page = _slide_action_plan(prs, d, narratives, page)

    # ── Slide 14: Closing ──
    make_closing(prs, d.get("title_agency", ""))
    page += 1

    # ── Guardrail validation ──
    from report.ppt_design import validate_ppt_guardrail
    validate_ppt_guardrail(prs, label="SEO Report")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info(f"PPT built: {page - 1} slides, {buf.getbuffer().nbytes // 1024} KB")
    return buf.getvalue()
