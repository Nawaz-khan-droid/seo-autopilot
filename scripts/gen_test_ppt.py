"""Generate a test PPT with sample data and save to output folder."""
import io, os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from report.facts import ReportFacts, RankingRow, Kpidata, TechnicalData, CWVData, MonthlyTrafficPoint
from report.evidence import Evidence
from report.ppt import build_ppt

f = ReportFacts()
f.metadata.client_name = 'KW Demo'
f.metadata.agency_name = 'SEO Agency'
f.metadata.report_month = 'June 2026'
f.metadata.generated_at = '2026-06-10T14:00:00'

kws = [
    ('seo services mumbai', '3', '+2'),
    ('digital marketing agency mumbai', '8', '-1'),
    ('website designing company mumbai', '12', '+5'),
    ('seo agency mumbai', '5', '0'),
    ('digital marketing company', '15', '-3'),
    ('best seo company mumbai', '1', '+1'),
    ('web development company mumbai', '20', '+2'),
    ('search engine optimization mumbai', '7', '0'),
    ('online marketing agency mumbai', '11', '+4'),
    ('seo expert mumbai', '6', '-2'),
]
for kw, pos, chg in kws:
    f.rankings.append(RankingRow(
        keyword=kw,
        target_url='https://example.com/' + kw.replace(' ', '-'),
        position=Evidence.verified(pos, 'SERP Snapshot'),
        change=Evidence.verified(chg, 'SERP History'),
        data_availability=Evidence.verified('Found', 'SERP Snapshot'),
    ))

f.kpis = Kpidata(
    clicks=Evidence.verified(1250, 'GSC API'),
    impressions=Evidence.verified(45200, 'GSC API'),
    monthly_traffic=[
        MonthlyTrafficPoint(month='Jan', users=Evidence.verified('2362','Analytics')),
        MonthlyTrafficPoint(month='Feb', users=Evidence.verified('2476','Analytics')),
        MonthlyTrafficPoint(month='Mar', users=Evidence.verified('4038','Analytics')),
        MonthlyTrafficPoint(month='Apr', users=Evidence.verified('4561','Analytics')),
        MonthlyTrafficPoint(month='May', users=Evidence.verified('5432','Analytics')),
        MonthlyTrafficPoint(month='Jun', users=Evidence.verified('4008','Analytics')),
    ],
)
f.technical = TechnicalData(
    health_score=Evidence.verified(72, 'Site Audit'),
    pages_audited=Evidence.verified(15, 'Site Audit'),
    total_issues=Evidence.verified(23, 'Site Audit'),
    missing_h1=Evidence.verified(3, 'Site Audit'),
    missing_meta=Evidence.verified(5, 'Site Audit'),
    missing_alt=Evidence.verified(42, 'Site Audit'),
    thin_pages=Evidence.verified(2, 'Site Audit'),
)
f.cwv = CWVData(
    mobile_score=Evidence.verified(65, 'PageSpeed API'),
    desktop_score=Evidence.verified(88, 'PageSpeed API'),
    lcp_seconds=Evidence.verified(3.2, 'PageSpeed API'),
    inp_ms=Evidence.verified(180, 'PageSpeed API'),
    cls_score=Evidence.verified(0.15, 'PageSpeed API'),
    render_blocking_ms=Evidence.verified(850, 'PageSpeed API'),
)

narrative = (
    'KW Demo tracked 10 keywords during June 2026. '
    'Of these, 5 improved, 2 dropped, and 3 remained stable. '
    'Top performer was best seo company mumbai moving to position 1. '
    'Mobile page speed needs attention at 65/100 vs desktop 88/100.'
)

ppt_bytes = build_ppt(f, narrative=narrative)
os.makedirs('output', exist_ok=True)
path = os.path.join('output', 'KW_Demo_June_2026_Report.pptx')
with open(path, 'wb') as fh:
    fh.write(ppt_bytes)
print(f'Saved: {path} ({len(ppt_bytes)//1024} KB)')

from pptx import Presentation
prs = Presentation(io.BytesIO(ppt_bytes))
print(f'Slides: {len(prs.slides)}')
for i, sl in enumerate(prs.slides, 1):
    txt = ''
    for sh in sl.shapes:
        if hasattr(sh, 'text') and sh.text.strip():
            txt = sh.text.strip()[:55]
            break
    print(f'  {i}: {txt}')
