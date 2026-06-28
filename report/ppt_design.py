"""Design system for agency-grade PPT reports.

Provides reusable shapes, tables, cards, headers, footers, and color constants.
All positions use a 12-column grid with 0.5" margins.

GUARDRAIL: All text sizes must stay above FONT_MIN (11pt).
Section headers must be >= FONT_SECTION (16pt).
Slide titles must be >= FONT_TITLE (20pt).
"""
from __future__ import annotations

import io
import logging
import re
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

# ── Font size constants ──
FONT_MIN     = 11   # Minimum allowed (guardrail enforced)
FONT_BODY    = 12   # Body / table text
FONT_SECTION = 16   # Section sub-headers
FONT_TITLE   = 20   # Slide titles (header_bar)
FONT_HEADING = 14   # Minor headings
FONT_CARD    = 9    # Card label/value text (sample uses 9pt)

# ── Slide content limits ──
MAX_PRIORITIES_PER_SLIDE = 7
MAX_ITEMS_PER_SEVERITY   = 3

# ── Forbidden patterns ──
FORBIDDEN_OVERFLOW_PATTERN = re.compile(r'\+\s*\d+\s*more', re.IGNORECASE)
FORBIDDEN_PRIORITY_P0      = re.compile(r'\bP0\b')

# ── Color Palette (sample template: #0F2747 navy, #2563EB blue) ──
NAVY          = RGBColor(0x0F, 0x27, 0x47)
NAVY_LIGHT    = RGBColor(0x2C, 0x3E, 0x66)
PRIMARY_BLUE  = RGBColor(0x25, 0x63, 0xEB)   # Accent blue (replaces TEAL)
DARK_BLUE     = RGBColor(0x1E, 0x40, 0xAF)   # Header badge bg
TEAL          = RGBColor(0x14, 0xB8, 0xA6)   # Secondary accent
AMBER         = RGBColor(0xF5, 0x9E, 0x0B)
RED           = RGBColor(0xDC, 0x26, 0x26)
GREEN         = RGBColor(0x16, 0xA3, 0x4A)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG       = RGBColor(0xF8, 0xFA, 0xFC)   # Very light gray
GRAY_LINE     = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_MID      = RGBColor(0x47, 0x55, 0x69)   # Card text
GRAY_TEXT     = RGBColor(0x47, 0x55, 0x69)   # Body text
DARK_TEXT     = RGBColor(0x0F, 0x17, 0x2A)   # Heading text
SUBTITLE_BLUE = RGBColor(0xDB, 0xEA, 0xFE)   # Header subtitle
INFO_BG       = RGBColor(0xEF, 0xF6, 0xFF)   # Blue-tinted section
GREEN_LIGHT   = RGBColor(0xEC, 0xFD, 0xF5)   # Green-tinted section
AMBER_LIGHT   = RGBColor(0xFF, 0xFB, 0xEB)   # Amber-tinted section
CORAL         = RGBColor(0xEF, 0x44, 0x44)   # Used by action_plan_doc

# ── Grid ──
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN  = 0.5
COL     = (SLIDE_W - 2 * MARGIN) / 12
CONTENT_W = SLIDE_W - 2 * MARGIN
CONTENT_L = MARGIN
CONTENT_R = SLIDE_W - MARGIN


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Reusable helpers ──

def _tb(slide, l, t, w, h, text="", size=12, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, name="Calibri", anchor=MSO_ANCHOR.TOP):
    """Add a textbox. Enforces FONT_MIN guardrail — silently clamps below it."""
    if size < FONT_MIN:
        logger.warning("_tb clamped font %dpt -> %dpt: '%s'", size, FONT_MIN, str(text)[:40])
        size = FONT_MIN
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tx.text_frame.word_wrap = True
    tx.text_frame.paragraphs[0].alignment = align
    p = tx.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    try:
        tx.text_frame.paragraphs[0].space_before = Pt(0)
        tx.text_frame.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass
    return tx


def _rect(slide, l, t, w, h, fill=NAVY):
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _line(slide, l, t, w, color=GRAY_LINE, width=1.0):
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(0.0)
    )
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


# ── Slide-level components ──

def add_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: str = "", page_num: int = 0,
                   section: str = ""):
    """Navy header (0.72\") + blue accent line (0.06\") + optional section badge."""
    _rect(slide, 0, 0, SLIDE_W, 0.72, NAVY)
    _rect(slide, 0, 0.72, SLIDE_W, 0.06, PRIMARY_BLUE)
    _tb(slide, MARGIN, 0.12, CONTENT_W - 1.5, 0.36, title,
        size=20, bold=True, color=WHITE)
    if subtitle:
        _tb(slide, MARGIN, 0.44, CONTENT_W - 2.0, 0.22, subtitle,
            size=9.5, color=SUBTITLE_BLUE)
    if section:
        badge_w = len(section) * 0.09 + 0.30
        _rect(slide, SLIDE_W - MARGIN - badge_w - 0.1, 0.18, badge_w + 0.2, 0.28, DARK_BLUE)
        _tb(slide, SLIDE_W - MARGIN - badge_w, 0.20, badge_w, 0.24, section,
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_footer(slide, agency: str = "", page_num: int = 0):
    """Footer: separator line + small text. No navy bar."""
    _line(slide, MARGIN, SLIDE_H - 0.5, CONTENT_W, GRAY_LINE, 0.5)
    _tb(slide, MARGIN, SLIDE_H - 0.47, 8.0, 0.20,
        agency or "SEO Report", size=8, color=GRAY_MID)
    if page_num:
        _tb(slide, SLIDE_W - MARGIN - 0.6, SLIDE_H - 0.47, 0.6, 0.20,
            str(page_num), size=8, color=GRAY_MID, align=PP_ALIGN.RIGHT)


# ── Content components ──

def kpi_card(slide, l, t, w, h, title, value, sub="",
             title_color=PRIMARY_BLUE, accent_color=None):
    """White card with left accent bar + small multi-line data."""
    if accent_color is None:
        accent_color = title_color
    _rect(slide, l, t, w, h, WHITE)
    _rect(slide, l, t, 0.07, h, accent_color)
    _tb(slide, l + 0.2, t + 0.08, w - 0.3, 0.25, title,
        size=9, bold=True, color=GRAY_MID)
    _tb(slide, l + 0.2, t + 0.30, w - 0.3, 0.30, str(value),
        size=14, bold=True, color=DARK_TEXT)
    if sub:
        _tb(slide, l + 0.2, t + 0.65, w - 0.3, 0.30, sub,
            size=9, color=GRAY_MID)


def metric_row(slide, l, t, w, label, value, status="neutral",
               badge_text=""):
    x = l
    _tb(slide, x, t, w * 0.35, 0.28, label, size=9, color=GRAY_TEXT)
    x += w * 0.35
    val_color = {"good": GREEN, "warn": AMBER, "bad": RED}.get(status, DARK_TEXT)
    _tb(slide, x, t, w * 0.25, 0.28, str(value), size=9, bold=True,
        color=val_color)
    x += w * 0.25
    if badge_text:
        badge_colors = {
            "verified": PRIMARY_BLUE, "observed": AMBER,
            "estimated": GRAY_MID, "no_data": RED,
        }
        bc = badge_colors.get(badge_text.lower(), GRAY_MID)
        _tb(slide, x, t + 0.01, w * 0.35, 0.24,
            f"[{badge_text}]", size=9, bold=True, color=bc)


# ── Tables ──

# ── Visual widgets ──

def progress_bar(slide, l, t, w, h, pct, bar_color=PRIMARY_BLUE, bg_color=GRAY_LINE):
    """Horizontal progress bar filled to pct (0-100)."""
    _rect(slide, l, t, w, h, bg_color)
    if pct > 0:
        fill_w = max(w * min(pct / 100.0, 1.0), 0.02)
        _rect(slide, l, t, fill_w, h, bar_color)


def big_stat(slide, l, t, w, h, number, label, color=PRIMARY_BLUE):
    """White card with left accent bar + number + label."""
    _rect(slide, l, t, w, h, WHITE)
    _rect(slide, l, t, 0.07, h, color)
    _tb(slide, l + 0.15, t + 0.08, w - 0.25, 0.35, str(number),
        size=18, bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    _tb(slide, l + 0.15, t + 0.42, w - 0.25, 0.25, label,
        size=9, color=GRAY_MID, align=PP_ALIGN.LEFT)


def status_badge(slide, l, t, text, color=PRIMARY_BLUE, size=9):
    """Colored badge/pill for status labels."""
    tw = len(text) * 0.09 + 0.3
    _rect(slide, l, t, tw, 0.24, color)
    _tb(slide, l + 0.04, t + 0.02, tw - 0.08, 0.20, text,
        size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Tables ──

def add_table(slide, l, t, col_widths, headers, rows,
              header_bg=NAVY, header_fg=WHITE, font_size=12):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    row_h = 0.35
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(l), Inches(t),
        Inches(total_w), Inches(row_h * n_rows)
    )
    table = table_shape.table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = Inches(cw)
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(hdr)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = header_fg
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for ri, row in enumerate(rows):
        bg = GRAY_BG if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val) if val is not None else "\u2014"
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


# ── Cover slide ──

def make_cover(prs: Presentation, client: str, month: str, agency: str,
               generated: str = "", health_score: int | None = None) -> Any:
    """Clean white cover with title, stat cards, and narrative panel (sample style)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    _ = agency  # unused on cover per sample style

    # Title
    _tb(slide, 0.55, 0.95, 3.5, 0.7,
        f"Monthly SEO Report\n{client} | {month}",
        size=28, bold=True, color=DARK_TEXT)

    # ── Left column: info cards ──
    score_str = f"{health_score} / 100" if health_score is not None else "\u2014 / 100"
    # Grade card
    _rect(slide, 0.55, 3.0, 3.2, 1.25, WHITE)
    _rect(slide, 0.55, 3.0, 0.07, 1.25, PRIMARY_BLUE)
    _tb(slide, 0.75, 3.08, 2.8, 0.25, "SEO Health Score",
        size=9, bold=True, color=GRAY_MID)
    _tb(slide, 0.75, 3.30, 2.8, 0.4, score_str,
        size=14, bold=True, color=DARK_TEXT)

    # Top Winner card
    _rect(slide, 0.55, 4.40, 3.2, 1.25, WHITE)
    _rect(slide, 0.55, 4.40, 0.07, 1.25, GREEN)
    _tb(slide, 0.75, 4.48, 2.8, 0.25, "Top Winner",
        size=9, bold=True, color=GRAY_MID)
    _tb(slide, 0.75, 4.70, 2.8, 0.5,
        "search engine optimization company mumbai\n+6 positions to #2",
        size=9, bold=True, color=GRAY_MID)

    # Bottleneck card
    _rect(slide, 0.55, 5.80, 3.2, 1.0, WHITE)
    _rect(slide, 0.55, 5.80, 0.07, 1.0, AMBER)
    _tb(slide, 0.75, 5.88, 2.8, 0.25, "Bottleneck",
        size=9, bold=True, color=GRAY_MID)
    _tb(slide, 0.75, 6.10, 2.8, 0.4, "Page speed gap\nMobile 73 vs Desktop 97",
        size=9, bold=True, color=GRAY_MID)

    # ── Right panel: narrative ──
    _rect(slide, 4.20, 1.30, 8.0, 5.50, WHITE)   # white backdrop
    _tb(slide, 4.50, 1.70, 7.4, 0.5, "Executive Focus",
        size=17, bold=True, color=DARK_TEXT)
    # First text block
    _rect(slide, 4.45, 2.20, 7.5, 1.85, GRAY_BG)
    _tb(slide, 4.60, 2.30, 7.2, 1.65,
        "What this month means\n\nTraffic increased strongly, with organic users up 55.4% "
        "and engaged sessions growing 20.9%. Core keywords hold position well.",
        size=12, color=DARK_TEXT)
    # Second text block
    _rect(slide, 4.45, 4.25, 7.5, 1.80, GRAY_BG)
    _tb(slide, 4.60, 4.35, 7.2, 1.60,
        "Recommended focus next month\n\nPrioritize technical fixes for H1, meta descriptions, "
        "and image alt text. Continue strengthening content for page 2 keywords.",
        size=12, color=DARK_TEXT)

    _line(slide, 0.45, SLIDE_H - 0.5, CONTENT_W, GRAY_LINE, 0.5)
    _tb(slide, 0.45, SLIDE_H - 0.47, 8.0, 0.20,
        f"{client} | Monthly SEO Report", size=8, color=GRAY_MID)
    return slide


def score_gauge(slide, l, t, w, h, score, label="SEO Score", color=PRIMARY_BLUE, max_score=100):
    """Large score display with label, like Woorank's score widget."""
    _rect(slide, l, t, w, h, GRAY_BG)
    _rect(slide, l, t, w, 0.04, color)
    mid_x = l + w / 2
    score_str = f"{int(score)}" if score is not None else "\u2014"
    _tb(slide, l + 0.15, t + 0.1, w - 0.3, 0.45, score_str,
        size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    _tb(slide, l + 0.15, t + 0.6, w - 0.3, 0.25, f"/ {max_score}",
        size=14, color=GRAY_MID, align=PP_ALIGN.CENTER)
    _tb(slide, l + 0.15, t + 0.8, w - 0.3, 0.25, label,
        size=12, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    progress_bar(slide, l + 0.15, t + 1.05, w - 0.3, 0.15, score or 0, color)


def checklist_group(slide, l, t, w, items, title="", title_color=PRIMARY_BLUE, item_color=DARK_TEXT):
    """Checklist-style list group with optional title."""
    y = t
    if title:
        _tb(slide, l, y, w, 0.3, title, size=14, bold=True, color=title_color)
        y += 0.35
    for item in items:
        _tb(slide, l + 0.1, y, w - 0.1, 0.28,
            f"\u2610 {item}", size=12, color=item_color)
        y += 0.28


def make_closing(prs: Presentation, agency: str) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    _rect(slide, 0, 0, SLIDE_W, 0.72, NAVY)
    _rect(slide, 0, 0.72, SLIDE_W, 0.06, PRIMARY_BLUE)
    badge_w = 0.7
    _rect(slide, SLIDE_W - MARGIN - badge_w - 0.1, 0.18, badge_w + 0.2, 0.28, DARK_BLUE)
    _tb(slide, SLIDE_W - MARGIN - badge_w, 0.20, badge_w, 0.24, "End",
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, 0.85, 1.55, 11.6, 4.55, NAVY)
    _tb(slide, MARGIN, 2.8, CONTENT_W, 0.8, "Thank You",
        size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(slide, MARGIN, 3.8, CONTENT_W, 0.5,
        "Let's discuss next steps on our call.",
        size=14, color=SUBTITLE_BLUE, align=PP_ALIGN.CENTER)
    add_footer(slide, f"{agency}")
    return slide


# ═══════════════════════════════════════════════════════════════
# LAYOUT COLLISION HELPERS
# ═══════════════════════════════════════════════════════════════

def _shape_bounds_inches(shape) -> tuple[float, float, float, float]:
    """Return (left, top, right, bottom) in inches from EMU."""
    left = shape.left / 914400
    top = shape.top / 914400
    right = left + (shape.width / 914400)
    bottom = top + (shape.height / 914400)
    return (left, top, right, bottom)


def _bbox_overlap(a: tuple, b: tuple, min_xy: float = 0.05) -> bool:
    """True if two normalized (L,T,R,B) boxes overlap by at least min_xy in each axis."""
    ox = min(a[2], b[2]) - max(a[0], b[0])
    oy = min(a[3], b[3]) - max(a[1], b[1])
    return ox > min_xy and oy > min_xy


def _bbox_contains(outer: tuple, inner: tuple) -> bool:
    """True if *outer* completely contains *inner* (within 0.01" slop)."""
    return (outer[0] - 0.01 <= inner[0] and outer[1] - 0.01 <= inner[1]
            and outer[2] + 0.01 >= inner[2] and outer[3] + 0.01 >= inner[3])


# ═══════════════════════════════════════════════════════════════
# GUARDRAIL ENFORCER — runs after every PPT / action-plan build
# ═══════════════════════════════════════════════════════════════

HEADER_BOTTOM = 0.9        # shapes fully above this are header decoration
FOOTER_TOP    = SLIDE_H - 0.55   # shapes at/below this are footer (line at 7.0)


def validate_ppt_guardrail(
    prs: Presentation,
    label: str = "PPT",
    raise_on_violation: bool = False,
) -> list[str]:
    """Post-generation guardrail checker.

    Scans every slide for:
      - Font sizes below FONT_MIN (11pt)
      - "+X more" hidden overflow text
      - P0 priority labels
      - Layout collisions (text boxes overlapping, shapes past footer)

    Returns a list of violation strings (empty = clean).
    If raise_on_violation=True AND violations exist, raises RuntimeError.
    """
    violations: list[str] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # ── Check every text run for font size ──
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.font.size and p.font.size < Pt(FONT_MIN):
                        snippet = (p.text or "")[:35].replace("\n", " ")
                        violations.append(
                            f"[{label} Slide {slide_num}] Font {int(p.font.size.pt)}pt < {FONT_MIN}pt: "
                            f"'{snippet}'"
                        )
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            if p.font.size and p.font.size < Pt(FONT_MIN):
                                snippet = (p.text or "")[:35].replace("\n", " ")
                                violations.append(
                                    f"[{label} Slide {slide_num}] Table font {int(p.font.size.pt)}pt < {FONT_MIN}pt: "
                                    f"'{snippet}'"
                                )

        # ── Check for hidden overflow text ──
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if FORBIDDEN_OVERFLOW_PATTERN.search(text):
                    violations.append(
                        f"[{label} Slide {slide_num}] Hidden overflow found: "
                        f"'{text[:50]}'"
                    )

        # ── Check for P0 priority labels ──
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if FORBIDDEN_PRIORITY_P0.search(text):
                    violations.append(
                        f"[{label} Slide {slide_num}] P0 priority label found: "
                        f"'{text[:50]}'"
                    )

        # ── Layout collision & boundary check ──
        content_shapes: list[tuple] = []  # (shape, L, T, R, B)
        for shape in slide.shapes:
            l, t, r, b = _shape_bounds_inches(shape)
            w, h = r - l, b - t
            # Skip zero-area shapes (lines) and header / footer decoration
            if w < 0.01 or h < 0.01:
                continue
            if b <= HEADER_BOTTOM or t >= FOOTER_TOP:
                continue
            content_shapes.append((shape, l, t, r, b))

        # ── Boundary overflow: any shape crossing below FOOTER_TOP ──
        for shape, l, t, r, b in content_shapes:
            if b > FOOTER_TOP:
                # Skip full-width footer bars (intentional design on cover/closing)
                if r - l >= SLIDE_W - 0.2 and b >= SLIDE_H - 0.05:
                    continue
                snippet = shape.text_frame.text[:45] if shape.has_text_frame else ""
                violations.append(
                    f"[{label} Slide {slide_num}] Shape past footer "
                    f"({l:.1f}\",{t:.1f}\")\u2192({r:.1f}\",{b:.1f}\"): '{snippet}'"
                )

        # ── Text-box overlap (text should never overlap text) ──
        txt_shapes = [(s, *bb) for s, *bb in content_shapes if s.has_text_frame]
        for i in range(len(txt_shapes)):
            for j in range(i + 1, len(txt_shapes)):
                s1, l1, t1, r1, b1 = txt_shapes[i]
                s2, l2, t2, r2, b2 = txt_shapes[j]
                box1 = (l1, t1, r1, b1)
                box2 = (l2, t2, r2, b2)
                if not _bbox_overlap(box1, box2):
                    continue
                # Intentional nesting (e.g. label inside a card) — not a collision
                if _bbox_contains(box1, box2) or _bbox_contains(box2, box1):
                    continue
                t1_text = (s1.text_frame.text or "").strip()[:30]
                t2_text = (s2.text_frame.text or "").strip()[:30]
                # Both blank — not a visual collision
                if not t1_text or not t2_text:
                    continue
                violations.append(
                    f"[{label} Slide {slide_num}] Text overlap: "
                    f"'{t1_text}' vs '{t2_text}' "
                    f"@ ({l1:.1f}\",{t1:.1f}\")\u2192({r1:.1f}\",{b1:.1f}\") / "
                    f"({l2:.1f}\",{t2:.1f}\")\u2192({r2:.1f}\",{b2:.1f}\")"
                )

    # ── Logging ──
    if violations:
        logger.warning("GUARDRAIL %d violation(s) in %s:", len(violations), label)
        for v in violations:
            logger.warning("  %s", v)
    else:
        logger.info("GUARDRAIL %s passed all checks.", label)

    if raise_on_violation and violations:
        raise RuntimeError(
            f"Guardrail failed with {len(violations)} violation(s) in {label}:\n"
            + "\n".join(violations[:10])
        )
    return violations
